#!/usr/bin/env python3
# -*- coding: utf-8 -*-

# The _XMLHandler class is:
# Copyright (c) Martin Blech
# Released under the MIT license.
# see https://opensource.org/licenses/MIT

__author__  = 'm.yama'
__license__ = 'MIT'
__date__    = 'Wed Jul 31 17:35:47 2019'
__version__ = '0.0.1'

__all__ = [
    "reader",
    "reader_groupby",
    "readerclip",
    "readerclip_groupby",
]

import os
import sys
import pathlib
import csv
import codecs
from io import IOBase, StringIO, BytesIO
from tarfile import ExFileObject, TarFile
from zipfile import ZipFile, ZipExtFile
from gzip import GzipFile
from bz2 import BZ2File
from lzma import LZMAFile
from rarfile import RarFile
from lhafile import LhaFile
import json as jsonlib
from html.parser import HTMLParser
from xml.parsers.expat import ParserCreate

from collections import namedtuple

#3rd party
try:
    from pptx import Presentation
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install python-pptx\n")
    Presentation = ImportError

try:
    from docx import Document
    from docx.oxml.text.run import CT_Text, CT_Br
    from docx.oxml.section import CT_PageMar
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install python-docx\n")
    Document = ImportError

try:
    import xlrd
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install xlrd\n")
    xlrd = ImportError

try:
    from pdfminer.pdfparser import PDFParser, PDFDocument
    from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
    from pdfminer.converter import PDFPageAggregator
    from pdfminer.layout import LAParams, LTTextBox, LTTextLine
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install pdfminer3k\n")
    PDFParser = ImportError
    PDFDocument = ImportError
    PDFResourceManager = ImportError
    PDFPageInterpreter = ImportError
    PDFPageAggregator = ImportError
    LAParams = ImportError
    LTTextBox = ImportError
    LTTextLine = ImportError

try:
    from pyperclip import paste as getclip
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install pyperclip\n")
    Presentation = ImportError

pinfo = namedtuple("LazyReader", ["path", "target", "value"])

def getencoding(dat:bytes):
    import nkf
    if b"\0" in dat:
        return None
    enc = nkf.guess(dat).lower()
    if enc and enc == "shift_jis":
        return "cp932"
    elif enc == "binary":
        return None
    else:
        return enc

def binopen(f, mode="rb", *args, **kw):

    check = lambda *tp: isinstance(f, tp)

    if check(str):
        return open(pathlib.Path(f), mode, *args, **kw)

    if check(pathlib.Path):
        return open(f, mode, *args, **kw)

    if check(ZipExtFile, ExFileObject, BytesIO):
        return f

    if check(bytearray, bytes):
        bio = BytesIO(f)
        bio.name = None
        return bio

    if check(IOBase):
        if check(StringIO):
            e = f.encoding
            if e:
                bio = BytesIO(f.getvalue().encode(e))
            else:
                bio = BytesIO(f.getvalue().encode())
            bio.name = f.name if hasattr(f, "name") else None
            return bio

        try:
            m = f.mode
        except AttributeError:
            m = f._mode
        if isinstance(m, int):
            return f

        name = f.name

        if "b" in m:
            return f
        else:
            return open(name, mode=m + "b")

    raise ValueError("Unknown Object `{}`. filename or filepointer buffer".format(type(f)))

def pathbin(path_or_buffer):
    fp = binopen(path_or_buffer)
    try:
        path = pathlib.Path(fp.name)
        return path, fp
    except AttributeError:
        return fp._fp.name, fp

def pptx(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    for i, s in enumerate(Presentation(fp).slides):
        for t in (r.text for sp in s.shapes if sp.has_text_frame for p in sp.text_frame.paragraphs for r in p.runs if r.text):
            yield pinfo(path, i, t)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def pptx_groupby(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    for i, s in enumerate(Presentation(fp).slides):
        yield pinfo(path, i, [sp.text for sp in s.shapes if sp.has_text_frame])

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def docx(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    i = 0
    it = Document(fp).element.getiterator()
    for txt in it:
        if isinstance(txt, CT_Br):
            i += 1
        if isinstance(txt, CT_Text):
            yield pinfo(path, i, txt.text)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def docx_groupby(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    i = 0
    it = Document(fp).element.getiterator()
    ret = []
    for txt in it:
        if isinstance(txt, (CT_Br, CT_PageMar)):
            yield pinfo(path, i, ret)
            ret = []
            i += 1
        if isinstance(txt, CT_Text):
            ret.append(txt.text)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def xlsx(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    with xlrd.open_workbook(file_contents=fp.read()) as wb:
        for sh in wb.sheets():
            sname = sh.name
            for row in sh.get_rows():
                yield pinfo(path, sname, row)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def xlsx_groupby(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    with xlrd.open_workbook(file_contents=fp.read()) as wb:
        for sh in wb.sheets():
            yield pinfo(path, sh.name, list(sh.get_rows()))

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def pdf(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    ps = PDFParser(fp)
    doc = PDFDocument()
    ps.set_document(doc)
    doc.set_parser(ps)
    doc.initialize('')

    mgr = PDFResourceManager(caching=True)
    dev = PDFPageAggregator(mgr, laparams=LAParams())
    ip = PDFPageInterpreter(mgr, dev)

    for i, page in enumerate(doc.get_pages(), 1):
        ip.process_page(page)
        text = "".join(x.get_text() for x in dev.get_result() if isinstance(x, (LTTextBox, LTTextLine)))
        yield pinfo(path, i, text.rstrip("{}\n".format(i)))

    dev.close()
    if not hasattr(path_or_buffer, "close"):
        fp.close()

def pdf_groupby(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    ps = PDFParser(fp)
    doc = PDFDocument()
    ps.set_document(doc)
    doc.set_parser(ps)
    doc.initialize('')

    mgr = PDFResourceManager(caching=True)
    dev = PDFPageAggregator(mgr, laparams=LAParams())
    ip = PDFPageInterpreter(mgr, dev)

    for i, page in enumerate(doc.get_pages(), 1):
        ip.process_page(page)
        rows = [x.get_text() for x in dev.get_result() if isinstance(x, (LTTextBox, LTTextLine))]
        yield pinfo(path, i, rows)

    dev.close()
    if not hasattr(path_or_buffer, "close"):
        fp.close()


def _iterdecoder(fp, encoding=None, buffer=1024**2):
    if encoding:
        while True:
            buf = fp.read(buffer) + fp.readline()
            if not buf:
                break
            for line in buf.decode(encoding).splitlines():
                yield line
    else:
        while True:
            buf = fp.read(buffer) + fp.readline()
            if not buf:
                break
            for line in buf.decode().splitlines():
                yield line

def xsv(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    dat = fp.read(1024) + fp.readline()

    e = getencoding(dat)
    txt = dat.decode(e)
    dialect = csv.Sniffer().sniff(txt)

    for row in csv.reader(StringIO(txt), dialect=dialect):
        yield pinfo(path, None, row)

    it = _iterdecoder(fp, e)
    for row in csv.reader(it, dialect=dialect):
        yield pinfo(path, None, row)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def xsv_groupby(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    dat = fp.read(1024) + fp.readline()

    e = getencoding(dat)
    txt = dat.decode(e)
    dialect = csv.Sniffer().sniff(txt)

    tmp = list(csv.reader(StringIO(txt), dialect=dialect))

    it = _iterdecoder(fp, e)
    yield pinfo(path, None, tmp + list(csv.reader(it, dialect=dialect)))

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def plantext(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    dat = fp.read(1024)
    fp.close()
    encoding = getencoding(dat)
    with codecs.open(path, encoding=encoding) as f:
        for line in f:
            yield pinfo(path, None, line)
    if not hasattr(path_or_buffer, "close"):
        fp.close()

def plantext_groupby(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    dat = fp.read(1024)
    fp.close()
    encoding = getencoding(dat)
    with codecs.open(path, encoding=encoding) as f:
        yield pinfo(path, None, f.readlines())

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def Zip(path_or_buffer, targets=[]):
    path, fp = pathbin(path_or_buffer)

    with ZipFile(fp) as f:
        if not targets:
            targets = f.namelist()
        for info in f.infolist():
            if info.is_dir():
                continue
            target = info.filename.encode("cp437").decode("cp932")
            if target in targets:
                dat = f.read(info)
                func = getreader(dat[:20], target)
                bio = BytesIO(dat)
                bio.name = target
                for row in func(bio):
                    yield pinfo(path, target, row.value)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def Zip_groupby(path_or_buffer, targets=[]):
    path, fp = pathbin(path_or_buffer)

    with ZipFile(path) as f:
        if not targets:
            targets = f.namelist()
        for info in f.infolist():
            if info.is_dir():
                continue
            target = info.filename.encode("cp437").decode("cp932")
            if target in targets:
                dat = f.read(info)
                func = getreader_groupby(dat[:20], target)
                bio = BytesIO(dat)
                bio.name = target
                for row in func(bio):
                    yield pinfo(path, target, row.value)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def gzip(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    target = path.stem
    with GzipFile(fileobj=fp) as f:
        func = getreader(f.read(20), target)
        f.seek(0)
        for row in func(f):
            yield pinfo(path, target, row)
    if not hasattr(path_or_buffer, "close"):
        fp.close()

def gzip_groupby(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    target = path.stem
    with GzipFile(fileobj=fp) as f:
        func = getreader_groupby(f.read(20), target)
        f.seek(0)
        for row in func(f):
            yield pinfo(path, target, row)
    if not hasattr(path_or_buffer, "close"):
        fp.close()

def bzip2(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    target = path.stem
    with BZ2File(fp) as f:
        func = getreader(f.read(20), target)
        f.seek(0)
        for row in func(f):
            yield pinfo(path, target, row)
    if not hasattr(path_or_buffer, "close"):
        fp.close()

def bzip2_groupby(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    target = path.stem
    with BZ2File(fp) as f:
        func = getreader_groupby(f.read(20), target)
        f.seek(0)
        for row in func(f):
            yield pinfo(path, target, row)
    if not hasattr(path_or_buffer, "close"):
        fp.close()

def lzma(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    target = path.stem
    with LZMAFile(fp) as f:
        func = getreader(f.read(20), target)
        f.seek(0)
        for row in func(f):
            yield pinfo(path, target, row)
    if not hasattr(path_or_buffer, "close"):
        fp.close()

def lzma_groupby(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    target = path.stem
    with LZMAFile(fp) as f:
        func = getreader_groupby(f.read(20), target)
        f.seek(0)
        for row in func(f):
            yield pinfo(path, target, row)
    if not hasattr(path_or_buffer, "close"):
        fp.close()

def tar(path_or_buffer, mode="r", targets=[]):
    path, fp = pathbin(path_or_buffer)

    with TarFile.open(mode=mode, fileobj=fp) as f:
        if not targets:
            targets = f.getnames()
        for info in f.getmembers():
            if info.isdir():
                continue
            target = info.name

            if target in targets:
                dat = f.extractfile(info).read()
                func = getreader(dat[:20], target)
                bio = BytesIO(dat)
                bio.name = target
                for row in func(bio):
                    yield pinfo(path, target, row.value)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def tar_groupby(path_or_buffer, mode="r", targets=[]):
    path, fp = pathbin(path_or_buffer)

    with TarFile.open(mode=mode, fileobj=fp) as f:
        if not targets:
            targets = f.getnames()
        for info in f.getmembers():
            if info.isdir():
                continue
            target = info.name

            if target in targets:
                dat = f.extractfile(info).read()
                func = getreader_groupby(dat[:20], target)
                bio = BytesIO(dat)
                bio.name = target
                for row in func(bio):
                    yield pinfo(path, target, row.value)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def lha(path_or_buffer, targets=[]):
    path, fp = pathbin(path_or_buffer)

    f = LhaFile(fp)
    if not targets:
        targets = f.namelist()
    for info in f.infolist():
        if info.is_dir():
            continue
        target = info.filename
        if target in targets:

            dat = f.read(info.filename)
            func = getreader(dat[:20], target)
            bio = BytesIO(dat)
            bio.name = target
            for row in func(bio):
                yield pinfo(path, target, row.value)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def lha_groupby(path_or_buffer, targets=[]):
    path, fp = pathbin(path_or_buffer)

    f = LhaFile(fp)
    if not targets:
        targets = f.namelist()
    for info in f.infolist():
        if info.is_dir():
            continue
        target = info.filename
        if target in targets:
            dat = f.read(info.filename)
            func = getreader_groupby(dat[:20], target)
            bio = BytesIO(dat)
            bio.name = target
            for row in func(bio):
                yield pinfo(path, target, row.value)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def rar(path_or_buffer, targets=[]):
    path, fp = pathbin(path_or_buffer)

    with RarFile(fp) as f:
        if not targets:
            targets = f.namelist()
        for info in f.infolist():
            if info.is_dir():
                continue
            target = info.filename
            if target in targets:
                dat = f.read(info)
                func = getreader(dat[:20], target)
                bio = BytesIO(dat)
                bio.name = target
                for row in func(bio):
                    yield pinfo(path, target, row.value)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def rar_groupby(path_or_buffer, targets=[]):
    path, fp = pathbin(path_or_buffer)

    with RarFile(fp) as f:
        if not targets:
            targets = f.namelist()
        for info in f.infolist():
            if info.is_dir():
                continue
            target = info.filename
            if target in targets:
                dat = f.read(info)
                func = getreader_groupby(dat[:20], target)
                bio = BytesIO(dat)
                bio.name = target
                for row in func(bio):
                    yield pinfo(path, target, row.value)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def json(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    yield pinfo(path, None, jsonlib.loads(fp.read()))

    if not hasattr(path_or_buffer, "close"):
        fp.close()

json_groupby = json

class Parser(HTMLParser):
    def __init__(self, data, encoding=None):
        super().__init__()
        self.result = []
        self.feed(data, encoding)

    def handle_data(self, data):
        self.result.append(data)

    def feed(self, data, encoding=None):
        if isinstance(data, bytes):
            if encoding:
                data = data.decode(encoding)
            else:
                data = data.decode()
        super().feed(data)

def html(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    dat = fp.read()
    if not hasattr(path_or_buffer, "close"):
        fp.close()
    p = Parser(dat, getencoding(dat))

    for r in p.result:
        yield pinfo(path, None, r)

    del p


def html_groupby(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    dat = fp.read()
    if not hasattr(path_or_buffer, "close"):
        fp.close()
    p = Parser(dat, getencoding(dat))
    yield pinfo(path, None, p.result)

    del p

class _XMLHandler(object):
    def __init__(self,
                 item_depth=0,
                 item_callback=lambda *args: True,
                 xml_attribs=True,
                 attr_prefix='@',
                 cdata_key='#text',
                 force_cdata=False,
                 cdata_separator='',
                 postprocessor=None,
                 dict_constructor=dict,
                 strip_whitespace=True,
                 nssep=':',
                 namespaces=None,
                 force_list=None,
                 comment_key='#comment'):
        self.path = []
        self.stack = []
        self.data = []
        self.item = None
        self.item_depth = item_depth
        self.xml_attribs = xml_attribs
        self.item_callback = item_callback
        self.attr_prefix = attr_prefix
        self.cdata_key = cdata_key
        self.force_cdata = force_cdata
        self.cdata_separator = cdata_separator
        self.postprocessor = postprocessor
        self.dict_constructor = dict_constructor
        self.strip_whitespace = strip_whitespace
        self.nssep = nssep
        self.namespaces = namespaces
        self.namespace_declarations = dict()
        self.force_list = force_list
        self.comment_key = comment_key

    def _build_name(self, full_name):
        if self.namespaces is None:
            return full_name
        i = full_name.rfind(self.nssep)
        if i == -1:
            return full_name
        namespace, name = full_name[:i], full_name[i+1:]
        try:
            short_namespace = self.namespaces[namespace]
        except KeyError:
            short_namespace = namespace
        if not short_namespace:
            return name
        else:
            return self.nssep.join((short_namespace, name))

    def _attrs_to_dict(self, attrs):
        if isinstance(attrs, dict):
            return attrs
        return self.dict_constructor(zip(attrs[0::2], attrs[1::2]))

    def startNamespaceDecl(self, prefix, uri):
        self.namespace_declarations[prefix or ''] = uri

    def startElement(self, full_name, attrs):
        name = self._build_name(full_name)
        attrs = self._attrs_to_dict(attrs)
        if attrs and self.namespace_declarations:
            attrs['xmlns'] = self.namespace_declarations
            self.namespace_declarations = dict()
        self.path.append((name, attrs or None))
        if len(self.path) > self.item_depth:
            self.stack.append((self.item, self.data))
            if self.xml_attribs:
                attr_entries = []
                for key, value in attrs.items():
                    key = self.attr_prefix+self._build_name(key)
                    if self.postprocessor:
                        entry = self.postprocessor(self.path, key, value)
                    else:
                        entry = (key, value)
                    if entry:
                        attr_entries.append(entry)
                attrs = self.dict_constructor(attr_entries)
            else:
                attrs = None
            self.item = attrs or None
            self.data = []

    def endElement(self, full_name):
        name = self._build_name(full_name)
        if len(self.path) == self.item_depth:
            item = self.item
            if item is None:
                item = (None if not self.data
                        else self.cdata_separator.join(self.data))

            should_continue = self.item_callback(self.path, item)
            if not should_continue:
                raise InterruptedError
        if len(self.stack):
            data = (None if not self.data
                    else self.cdata_separator.join(self.data))
            item = self.item
            self.item, self.data = self.stack.pop()
            if self.strip_whitespace and data:
                data = data.strip() or None
            if data and self.force_cdata and item is None:
                item = self.dict_constructor()
            if item is not None:
                if data:
                    self.push_data(item, self.cdata_key, data)
                self.item = self.push_data(self.item, name, item)
            else:
                self.item = self.push_data(self.item, name, data)
        else:
            self.item = None
            self.data = []
        self.path.pop()

    def characters(self, data):
        if not self.data:
            self.data = [data]
        else:
            self.data.append(data)

    def comments(self, data):
        if self.strip_whitespace:
            data = data.strip()
        self.item = self.push_data(self.item, self.comment_key, data)

    def push_data(self, item, key, data):
        if self.postprocessor is not None:
            result = self.postprocessor(self.path, key, data)
            if result is None:
                return item
            key, data = result
        if item is None:
            item = self.dict_constructor()
        try:
            value = item[key]
            if isinstance(value, list):
                value.append(data)
            else:
                item[key] = [value, data]
        except KeyError:
            if self._should_force_list(key, data):
                item[key] = [data]
            else:
                item[key] = data
        return item

    def _should_force_list(self, key, value):
        if not self.force_list:
            return False
        if isinstance(self.force_list, bool):
            return self.force_list
        try:
            return key in self.force_list
        except TypeError:
            return self.force_list(self.path[:-1], key, value)


def xmlparse(data:bytes, encoding=None, namespaces=False, nssep=':', comments=False, **kw):
    hl = _XMLHandler(nssep=nssep, **kw)
    ps = ParserCreate(encoding, nssep or None)
    ps.ordered_attributes = True
    ps.StartNamespaceDeclHandler = hl.startNamespaceDecl
    ps.StartElementHandler = hl.startElement
    ps.EndElementHandler = hl.endElement
    ps.CharacterDataHandler = hl.characters

    if comments:
        ps.CommentHandler = hl.comments

    ps.buffer_text = True
    ps.Parse(data, True)

    return hl.item

def xml(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    dat = fp.read()
    if not hasattr(path_or_buffer, "close"):
        fp.close()

    e = getencoding(dat)
    yield pinfo(path, None, xmlparse(dat, e))

xml_groupby = xml

def mdb(path_or_buffer):
    path, fp = pathbin(path_or_buffer)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def getreader(dat:bytes, path:str):
    ext = os.path.splitext(path)[1].lower()
    startsdat = dat.startswith
    startsext = ext.startswith
    func = None

    if b"\x00" in dat:
        if startsdat(b"PK"):
            handle_o = (
                (".ppt", pptx),
                (".doc", docx),
                (".xls", xlsx),
            )
            func = next((fc for x, fc in handle_o if startsext(x)), Zip)
        else:
            handle_tz = (
                (b"\x1f\x8b", gzip, "r:gz"),
                (b"BZh", bzip2, "r:bz2"),
                (b"\xFD7zXZ\x00", lzma, "r:xz"),
            )

            func_mode = next((fc for x, fc, mode in handle_tz if startsdat(x)), None)
            if func_mode:
                func, mode = func_mode
                if b".tar\x00" in dat.lower() or b".tgz\x00" in dat.lower():
                    func = lambda x: tar(x, mode)

            else:
                handle_d = (
                    (b"%PDF-", pdf),
                    (b"ustar\x00", tar),
                    (b"ustar\x40", tar),
                    (b"ustar  \x00", tar),
                    (b"ustar  \x40", tar),
                    (b"!\xd1-lh0-", lha),
                    (b"!\xd1-lh1-", lha),
                    (b"!", lha),
                    (b"!\xd1-lh5-", lha),
                    (b"!\xd1-lh6-", lha),
                    (b"!\xd1-lh7-", lha),
                    (b"!\x82-lh0-", lha),
                    (b"!\x82-lh1-", lha),
                    (b"!\x82-lh4-", lha),
                    (b"!\x82-lh5-", lha),
                    (b"!\x82-lh6-", lha),
                    (b"!\x82-lh7-", lha),
                    (b"7zBCAF271C", lzma),
                    (b"RE\x7e\x5e", rar),
                    (b"Rar!\x1A\x07\x00", rar),
                    (b"Rar!\x1A\x07\x01\x00(RAR5)", rar),
                    (b"LZIP", lzma),
                    (b"\x00\x01\x00\x00Standard Jet DB\x00", mdb), #TODO
                    (b'\x00\x01\x00\x00Standard ACE DB\x00', mdb), #TODO
                    (b'SQLite format 3', None), #TODO
                )

                func = next((fc for x, fc in handle_d if startsdat(x)), None)
    else:
        handle_e = (
            (".csv", xsv),
            (".tsv", xsv),
            (".json", json),
            (".htm", html),
            (".xml", xml),
            ("txt", plantext),
            )
        func = next((fc for x, fc in handle_e if startsext(x)), None)
    if not func:
        raise ValueError("Unknown data type " + str(path))
    return func

def getreader_groupby(dat:bytes, path:str):
    ret = getreader(dat, path)
    return eval(ret.__name__ + "_groupby")


def reader(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    func = getreader(fp.read(20), path)
    fp.seek(0)
    return func(fp)

def reader_groupby(path_or_buffer):
    path, fp = pathbin(path_or_buffer)
    func = getreader_groupby(fp.read(20), path)
    fp.seek(0)
    return func(fp)

def readerclip(ftype="xsv"):
    """
    clipboard read parse

    ftype : xsv or html, xml, plantext, json
    """
    funcs = globals().copy()
    if ftype not in funcs:
        raise AttributeError

    cdat = getclip()
    with BytesIO() as bio:
        try:
            bdat = cdat.encode("utf-8")
        except:
            bdat = cdat.encode()
        bio.write(bdat)
        bio.seek(0)
        bio.name = "clipboad"
        for r in eval(ftype)(bio):
            yield r

def readerclip_groupby(ftype="xsv"):
    """
    clipboard read parse

    ftype : xsv or html, xml, plantext, json
    """
    funcs = globals().copy()
    if ftype not in funcs:
        raise AttributeError

    cdat = getclip()
    with BytesIO() as bio:
        try:
            bdat = cdat.encode("utf-8")
        except:
            bdat = cdat.encode()
        bio.write(bdat)
        bio.seek(0)
        bio.name = "clipboad"
        return eval(ftype + "_groupby")(bio)

def test():
    from datetime import datetime as dt
    from util.core import tdir

    def test_binopen():
        f = binopen(tdir+"test.csv")
        binopen(f)
        binopen(pathlib.Path(tdir+"test.csv"))

    def test_getencoding():
        with open(tdir+"test.csv", "rb") as f:
            e = getencoding(f.read())
            assert(e == "cp932")

    def test_docx():
        assert(len(list(docx(tdir+"test.docx"))) == 7)

    def test_docx_groupby():
        assert(len(list(docx_groupby(tdir+"test.docx"))) == 2)

    def test_xlsx():
        list(xlsx(tdir+"diff1.xlsx"))

    def test_xsv():
        assert(len(list(xsv(tdir+"test.csv"))) == 3)

    def test_gzip():
        list(gzip(tdir+"test.csv.gz"))
        list(gzip_groupby(tdir+"test.csv.gz"))

    def test_bzip2():
        list(bzip2(tdir+"test.csv.bz2"))
        list(bzip2_groupby(tdir+"test.csv.bz2"))

    def test_lzma():
        list(lzma(tdir+"test.csv.xz"))
        list(lzma_groupby(tdir+"test.csv.xz"))

    def test_zip():
        list(Zip_groupby(tdir+"test.zip"))

    def test_tar():
        list(tar(tdir+"test.tar"))
        list(tar_groupby(tdir+"test.tar"))

    def test_rarfile():
        list(rar(tdir+"test.rar"))
        list(rar_groupby(tdir+"test.rar"))

    def test_lha():
        list(lha(tdir+"test.lzh"))
        list(lha_groupby(tdir+"test.lzh"))

    def test_json():
        list(json(tdir+"test.json"))
        list(json_groupby(tdir+"test.json"))

    def test_html():
        list(html(tdir+"test.html"))
        list(html_groupby(tdir+"test.html"))

    def test_xml():
        list(xml(tdir+"test.xml"))
        list(xml_groupby(tdir+"test.xml"))

    def test_reader():
        list(reader(tdir+"diff1.xlsx"))

    for x, func in list(locals().items()):
        if x.startswith("test_") and callable(func):
            t1 = dt.now()
            func()
            t2 = dt.now()
            print("{} : time {}".format(x, t2-t1))

if __name__ == "__main__":
    test()
