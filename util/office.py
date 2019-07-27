#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""

My very usefull tools library
require pandas and dask!!

MIT License

"""

__all__ = [
    "reader",
    "readlines",
    "iterlines",
]


from collections import namedtuple
from pathlib import Path
import sys, os
import csv
import codecs
from io import IOBase, StringIO, BytesIO
from tarfile import ExFileObject
from zipfile import ZipExtFile

BASE_TYPE = [int, float, str, bytes, bytearray, bool]

def flatten(x):
    if x is None or type(x) in BASE_TYPE:
        return x
    return [z for y in x for z in ([y] if y is None or type(y) in BASE_TYPE else flatten(y))]

#3rd party
try:
    from pptx import Presentation
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install python-pptx\n")
    Presentation = ImportError

try:
    from docx import Document
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install python-docx\n")
    Document = ImportError

try:
    import xlrd
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install xlrd\n")
    xlrd = ImportError

try:
    from pdfminer.pdfparser import PDFParser, PDFDocument
    from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
    from pdfminer.converter import PDFPageAggregator
    from pdfminer.layout import LAParams, LTTextBox, LTTextLine
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install pdfminer3k\n")
    PDFParser = ImportError
    PDFDocument = ImportError
    PDFResourceManager = ImportError
    PDFPageInterpreter = ImportError
    PDFPageAggregator = ImportError
    LAParams = ImportError
    LTTextBox = ImportError
    LTTextLine = ImportError

def binopen(f, mode="rb", *args, **kw):

    check = lambda *tp: isinstance(f, tp)

    if check(str):
        return open(Path(f), mode, *args, **kw)

    if check(Path):
        return open(f, mode, *args, **kw)

    if check(ExFileObject, ZipExtFile, BytesIO):
        return f

    if check(bytearray, bytes):
        bio = BytesIO(f)
        bio.name = None
        return bio

    if check(IOBase):
        if check(StringIO):
            e = f.encoding
            if e:
                bio = BytesIO(f.getvalue().encode(e))
            else:
                bio = BytesIO(f.getvalue().encode())
            bio.name = f.name if hasattr(f, "name") else None
            return bio

        m = f.mode
        if isinstance(m, int):
            return f

        name = f.name

        if "b" in m:
            return f
        else:
            return open(name, mode=m + "b")

    raise ValueError("Unknown Object `{}`. filename or filepointer buffer".format(type(f)))

def splitfileobj(path_or_buffer, mode="rb"):
    fp = binopen(path_or_buffer)
    path = Path(fp.name)
    
#    try:
#        path = Path(path_or_buffer)
#        fp = binopen(path_or_buffer)
#    except TypeError:
#        path = Path(path_or_buffer.name)
#        fp = binopen(path_or_buffer)
    return path, fp

def getencoding(dat:bytes):
    import nkf
    if b"\0" in dat:
        return None
    enc = nkf.guess(dat).lower()
    if enc and enc == "shift_jis":
        return "cp932"
    elif enc == "binary":
        return None
    else:
        return enc

pinfo = namedtuple("OfficeDoc", ["path", "target", "value"])
def pptx(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)

    for i, s in enumerate(Presentation(fp).slides):
        for t in (r.text for sp in s.shapes if sp.has_text_frame for p in sp.text_frame.paragraphs for r in p.runs if r.text):
            yield pinfo(path, i, t)

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def pptx_groupby(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)

    for i, s in enumerate(Presentation(fp).slides):
        yield pinfo(path, i, [sp.text for sp in s.shapes if sp.has_text_frame])

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def docx(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)

    return (pinfo(path, None, txt) for txt in (p.text for p in Document(fp).paragraphs if p.text))

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def docx_groupby(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)
    for p in Document(fp).paragraphs:
        txt = p.text
        if txt:
            yield pinfo(path, None, txt)
    
    if not hasattr(path_or_buffer, "close"):
        fp.close()

def xlsx(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)

    with xlrd.open_workbook(file_contents=fp.read()) as wb:
        for i, sh in ((r, sh) for sh in wb.sheets() for r in range(sh.nrows)):
            yield pinfo(path, sh.name, sh.row_values(i))

    if not hasattr(path_or_buffer, "close"):
        fp.close()

def xlsx_groupby(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)

    with xlrd.open_workbook(file_contents=fp.read()) as wb:
        for sh in wb.sheets():
            yield pinfo(path, sh.name, [sh.row_values(i) for i in range(sh.nrows)])
    
    if not hasattr(path_or_buffer, "close"):
        fp.close()

def pdf(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)

    ps = PDFParser(fp)
    doc = PDFDocument()
    ps.set_document(doc)
    doc.set_parser(ps)
    doc.initialize('')

    mgr = PDFResourceManager(caching=True)
    dev = PDFPageAggregator(mgr, laparams=LAParams())
    ip = PDFPageInterpreter(mgr, dev)

    for i, page in enumerate(doc.get_pages(), 1):
        ip.process_page(page)
        text = "".join(x.get_text() for x in dev.get_result() if isinstance(x, (LTTextBox, LTTextLine)))
        yield pinfo(path, i, text.rstrip("{}\n".format(i)))

    dev.close()
    if not hasattr(path_or_buffer, "close"):
        fp.close()

def pdf_groupby(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)

    ps = PDFParser(fp)
    doc = PDFDocument()
    ps.set_document(doc)
    doc.set_parser(ps)
    doc.initialize('')

    mgr = PDFResourceManager(caching=True)
    dev = PDFPageAggregator(mgr, laparams=LAParams())
    ip = PDFPageInterpreter(mgr, dev)

    for i, page in enumerate(doc.get_pages(), 1):
        ip.process_page(page)
        rows = [x.get_text() for x in dev.get_result() if isinstance(x, (LTTextBox, LTTextLine))]
        yield pinfo(path, i, rows)

    dev.close()
    if not hasattr(path_or_buffer, "close"):
        fp.close()
        
def xsv(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)
    dat = fp.read(1024)
    fp.close()
    encoding = getencoding(dat)
    dialect = csv.Sniffer().sniff(dat.decode(encoding))
    with codecs.open(path, encoding=encoding) as f:
        for row in csv.reader(f, dialect=dialect):
            yield pinfo(path, None, row)

def xsv_groupby(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)
    dat = fp.read(1024)
    fp.close()
    encoding = getencoding(dat)
    dialect = csv.Sniffer().sniff(dat.decode(encoding))
    with codecs.open(path, encoding=encoding) as f:
        yield pinfo(path, None, list(csv.reader(f, dialect=dialect)))

def plantext(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)
    dat = fp.read(1024)
    fp.close()
    encoding = getencoding(dat)
    with codecs.open(path, encoding=encoding) as f:
        for line in f:
            yield pinfo(path, None, line)

def plantext_groupby(path_or_buffer):
    path, fp = splitfileobj(path_or_buffer)
    dat = fp.read(1024)
    fp.close()
    encoding = getencoding(dat)
    with codecs.open(path, encoding=encoding) as f:
        yield pinfo(path, None, f.readlines())

comp_magics = [
    b"\x1f\x8b",
    b"ustar\x00",
    b"ustar\x40",
    b"ustar  \x00",
    b"ustar  \x40",
    b"PK\x03\x04",
    b"PK\x05\x06",
    b"PK\x07\x08",
    b"BZh",
    b"!\xd1-lh0-",
    b"!\xd1-lh1-",
    b"!\xd1-lh4-",
    b"!\xd1-lh5-",
    b"!\xd1-lh6-",
    b"!\xd1-lh7-",
    b"!\x82-lh0-",
    b"!\x82-lh1-",
    b"!\x82-lh4-",
    b"!\x82-lh5-",
    b"!\x82-lh6-",
    b"!\x82-lh7-",
    b"\x1f\x9d", #CompressZFile
    b"\xFD7zXZ\x00",
    b"7zBCAF271C",
    b"RE\x7e\x5e",
    b"Rar!\x1A\x07\x00",
    b"Rar!\x1A\x07\x01\x00RAR5)",
    b"LZIP",
    b"MSCF\x00\x00\x00\x00",
    b"DGCA",
    b"GCAX",
]

def is_compress(path_or_buffer):
    fp = binopen(path_or_buffer)
    isin = fp.read(15).startswith
    fp.close()
    return sum([1 for x in comp_magics if isin(x)]) > 0

def compress(payh_or_buffer):
    return NotImplemented#TODO
    
def compress_groupby(payh_or_buffer):
    return NotImplemented#TODO
    
def json(path_or_buffer):
    return NotImplemented#TODO
    
def json_groupby(path_or_buffer):
    return NotImplemented#TODO
    
def xml(path_or_buffer):
    return NotImplemented#TODO

def xml_groupby(path_or_buffer):
    return NotImplemented#TODO
    
def reader(path_or_buffer):
    if hasattr(path_or_buffer, "name"):
        ext = os.path.splitext(path_or_buffer.name)[1]
    else:
        ext = os.path.splitext(path_or_buffer)[1]

    extfull = ext.lower()
    ext = extfull[:4]

    if ext == ".ppt":
        return pptx(path_or_buffer)
    elif ext == ".doc":
        return docx(path_or_buffer)
    elif ext == ".xls":
        return xlsx(path_or_buffer)
    elif ext == ".pdf":
        return pdf(path_or_buffer)
    elif ext in [".csv", ".tsv"]:
        return xsv(path_or_buffer)
    elif extfull == ".json":
        return json(path_or_buffer)
    elif ext in [".htm", ".xml"]:
        return xml(path_or_buffer)
    elif is_compress(path_or_buffer):
        return compress(path_or_buffer)
    else:
        return plantext(path_or_buffer)

def reader_groupby(path_or_buffer):
    if hasattr(path_or_buffer, "name"):
        ext = os.path.splitext(path_or_buffer.name)[1]
    else:
        ext = os.path.splitext(path_or_buffer)[1]

    extfull = ext.lower()
    ext = extfull[:4]

    if ext == ".ppt":
        return pptx_groupby(path_or_buffer)
    elif ext == ".doc":
        return docx_groupby(path_or_buffer)
    elif ext == ".xls":
        return xlsx_groupby(path_or_buffer)
    elif ext == ".pdf":
        return pdf_groupby(path_or_buffer)
    elif ext in [".csv", ".tsv"]:
        return xsv_groupby(path_or_buffer)
    elif extfull == ".json":
        return json_groupby(path_or_buffer)
    elif ext in [".htm", ".xml"]:
        return xml_groupby(path_or_buffer)
    elif is_compress(path_or_buffer):
        return compress_groupby(path_or_buffer)
    else:
        return plantext_groupby(path_or_buffer)
        
#def readlines(path_or_buffer):
#    return [r.value for r in reader(path_or_buffer)]

#def iterlines(path_or_buffer):
#    return (r.value for r in reader(path_or_buffer))

#from util.core import tdir
#for x in pptx_groupby(tdir+"test.pptx"):
#    print(x.value)
    
