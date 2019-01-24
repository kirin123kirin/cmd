#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""

My very usefull tools library
require pandas and dask!!

MIT License

"""
from collections import namedtuple
import sys

#3rd party
try:
    from pptx import Presentation
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install python-pptx\n")
    Presentation = ImportError

try:
    from docx import Document
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install python-docx\n")
    Document = ImportError

try:
    import xlrd
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install xlrd\n")
    xlrd = ImportError

try:
    from pdfminer.pdfparser import PDFParser, PDFDocument
    from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
    from pdfminer.converter import PDFPageAggregator
    from pdfminer.layout import LAParams, LTTextBox, LTTextLine
except ModuleNotFoundError:
    sys.stderr.write("** No module warning **\nPlease Install command: pip3 install pdfminer3k\n")
    PDFParser = ImportError
    PDFDocument = ImportError
    PDFResourceManager = ImportError
    PDFPageInterpreter = ImportError
    PDFPageAggregator = ImportError
    LAParams = ImportError
    LTTextBox = ImportError
    LTTextLine = ImportError

#my library
from util.core import Path, command
from io import StringIO

pinfo = namedtuple("OfficeDoc", ["path", "target", "value"])
def pptx(ppt_path_or_buffer):
    path = Path(ppt_path_or_buffer)
    for i, s in enumerate(Presentation(path).slides):
        for txt in (r.text for sp in s.shapes if sp.has_text_frame for p in sp.text_frame.paragraphs for r in p.runs if r.text):
            yield pinfo(path, i, txt)

def docx(docx_path_or_buffer):
    path = Path(docx_path_or_buffer)
    return (pinfo(path, "#TODO", txt) for txt in (p.text for p in Document(path).paragraphs if p.text))

def xlsx(xlsx_path_or_buffer):
    path = Path(xlsx_path_or_buffer)
    with xlrd.open_workbook(path) as wb:
        for i, sh in ((r, sh) for sh in wb.sheets() for r in range(sh.nrows)):
            yield pinfo(path, "{}:{}".format(sh.name, i+1), sh.row_values(i))

def pdf(path_or_buffer):
    path = Path(path_or_buffer)

    with path.open('rb') as fp:
        ps = PDFParser(fp)
        doc = PDFDocument()
        ps.set_document(doc)
        doc.set_parser(ps)
        doc.initialize('')

        mgr = PDFResourceManager(caching=True)
        dev = PDFPageAggregator(mgr, laparams=LAParams())
        ip = PDFPageInterpreter(mgr, dev)

        for i, page in enumerate(doc.get_pages(), 1):
            ip.process_page(page)
            text = "".join(x.get_text() for x in dev.get_result() if isinstance(x, (LTTextBox, LTTextLine)))
            yield pinfo(path, i, text.rstrip("{}\n".format(i)))

        dev.close()

def txt(path_or_buffer):
    path = Path(path_or_buffer)
    with path.open() as r:
        for i, line in enumerate(r, 1):
            yield pinfo(path, i, line.rstrip())


_handler = {
    ".ppt" : pptx,
    ".doc" : docx,
    ".xls" : xlsx,
    ".pdf" : pdf,
    ".txt" : txt,
}
def any(path_or_buffer):
    path = Path(path_or_buffer)
    try:
        ext = path.ext.lower()[:4]
        func = _handler.get(ext, txt)
        return func(path_or_buffer)
    except:
        sio = StringIO(command("xdoc2txt -8 {}".format(path)))
        return (pinfo(path, "", line) for line in sio if line.strip())

def test():
    from util.core import tdir
    from datetime import datetime as dt

    def test_pptx():
        f = tdir+"test.pptx"
        for x in pptx(f):
            assert(type(x) == pinfo)

    def test_docx():
        f = tdir+"test.docx"
        for x in docx(f):
            assert(type(x) == pinfo)

    def test_xlsx():
        f = tdir+'diff1.xlsx'
        for x in xlsx(f):
            assert(type(x) == pinfo)

    def test_pdf():
        f = tdir + "test.pdf"
        for x in pdf(f):
            assert(type(x) == pinfo)

    def test_txt():
        f = tdir + "test.csv"
        for x in txt(f):
            assert(type(x) == pinfo)

    def test_any():
        f = tdir+"test.pptx"
        for x in any(f):
            assert(type(x) == pinfo)

        f = tdir+"test.docx"
        for x in any(f):
            assert(type(x) == pinfo)

        f = tdir+'diff1.xlsx'
        for x in any(f):
            assert(type(x) == pinfo)

        f = tdir + "test.pdf"
        for x in any(f):
            assert(type(x) == pinfo)

        f = tdir + "test.csv"
        for x in any(f):
            assert(type(x) == pinfo)

    for x, func in list(locals().items()):
        if x.startswith("test_") and callable(func):
            t1 = dt.now()
            func()
            t2 = dt.now()
            print("{} : time {}".format(x, t2-t1))

def main():
    for t in any(sys.argv[1]):
        print("\t".join([str(t.path), t.target, " ".join(map(str,t.value)) if hasattr(t.value, "__iter__") else t.value]))

if __name__ == "__main__":
    #test()

    main()