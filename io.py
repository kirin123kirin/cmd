#!/usr/bin/env python3
# -*- coding: utf-8 -*-

# The _XMLHandler class is:
# Copyright (c) Martin Blech
# Released under the MIT license.
# see https://opensource.org/licenses/MIT
__author__  = 'm.yama'
__license__ = 'MIT'
__date__    = 'Wed Jul 31 17:35:47 2019'
__version__ = '0.0.1'

__all__ = [
    "lsdir",
    "xmltodict",
    "readrow",
    "dumper",
    "grouprow",
    "getinfo",
    "getsize",
    "DBrow",
    "DBgrouprow",
    "Path",
    "unicode_escape",
]

SAMPLEBYTE = 1024

import os
import re
import sys
import pathlib
import csv as _csv
import codecs
from io import StringIO, BytesIO
from tarfile import TarFile
from zipfile import ZipFile
from gzip import GzipFile
from bz2 import BZ2File
from lzma import LZMAFile
from rarfile import RarFile
from lhafile import LhaFile
import json as jsonlib
from html.parser import HTMLParser
from xml.parsers.expat import ParserCreate
from datetime import datetime
try:
    from lxml import etree
    from lxml.etree import _Element as Element
except ModuleNotFoundError:
    from xml.etree import ElementTree as etree
    from xml.etree.ElementTree import Element

from collections import namedtuple
from functools import lru_cache
from urllib.parse import quote_plus
import struct
from pickle import load as pkload

class NotInstalledModuleError(Exception):
    def stderr(self):
        raise __class__("** {} **".format(*self.args)) if self.args else __class__
    def __call__(self, *args, **kw): self.stderr()
    def __getattr__(self, *args, **kw): self.stderr()

#3rd party

## pptx
try:
    from pptx import Presentation
except ModuleNotFoundError:
    Presentation = NotInstalledModuleError("Please Install command: pip3 install python-pptx")

## ppt(OLE)
try:
    from msodumper import globals as mgl
    from msodumper.pptstream import PPTFile
except ModuleNotFoundError:
    mgl = PPTFile = NotInstalledModuleError("Please git clone: pip3 install git+https://github.com/kirin123kirin/mso-dumper.git")

## xls/xlsx
try:
    import xlrd
except ModuleNotFoundError:
    xlrd = NotInstalledModuleError("Please Install command: pip3 install xlrd")

## docx
try:
    from docx import Document
    from docx.oxml.text.run import CT_Text, CT_Br
    from docx.oxml.section import CT_PageMar
except ModuleNotFoundError:
    Document = CT_Text = CT_Br = CT_PageMar = NotInstalledModuleError("Please Install command: pip3 install python-docx")

## doc(OLE)
try:
    import olefile
    class MSWordOLE:
        """
            Reference github.com/dayflower/msworddoc-extractor
        """

        OFFSET_FIB_CCP_MAP = dict(
            Text    = 0x004c,
            Ftn     = 0x0050,
            Hdd     = 0x0054,
            Mcr     = 0x0058,
            Atn     = 0x005c,
            Edn     = 0x0060,
            Txbx    = 0x0064,
            HdrTxbx = 0x0068,
        )

        re_formatbin = re.compile(b"([\x07]*)[\x07]{2}")

        def __init__(self, path_or_buffer):
            self.path_or_buffer = path_or_buffer
            self.flag = {}
            self.ccp = {}
            self._ole = None
            self._load()

        def __enter__(self):
            return self

        def __exit__(self, ex_type, ex_value, trace):
            self.close()

        def close(self):
            self._ole.close()
            self._ole = None

        @property
        def ole(self):
            if self._ole is None:
                with binopen(self.path_or_buffer) as fp:
                    self._ole = olefile.OleFileIO(fp.read())
            return self._ole

        def _load(self):
            self.parse_fib(self.ole.openstream('WordDocument'))
            name_of_table = "{}Table".format(1 if self.flag["fWhichTblStm"] else 0)
            self.parse_piece_table(self.ole.openstream(name_of_table))

        def whole_contents(self, *args):
            return self.retrieve_and_filter(0, -1, *args)

        def document(self, *args):
            return self.retrieve_and_filter(0, self.ccp["Text"], *args)

        def footnote(self, *args):
            return self.retrieve_and_filter(self.ccp["Text"], self.ccp["Ftn"], *args)

        def header(self, *args):
            skips = [ "Text", "Ftn" ]
            return self.retrieve_token_and_filter(skips, "Hdd", *args)

        def macro(self, *args):
            skips = [ "Text", "Ftn", "Hdd" ]
            return self.retrieve_token_and_filter(skips, "Mcr", *args)

        def annotation(self, *args):
            skips = [ "Text", "Ftn", "Hdd", "Mcr" ]
            return self.retrieve_token_and_filter(skips, "Atn", *args)

        def endnote(self, *args):
            skips = [ "Text", "Ftn", "Hdd", "Mcr", "Atn" ]
            return self.retrieve_token_and_filter(skips, "Edn", *args)

        def textbox(self, *args):
            skips = [ "Text", "Ftn", "Hdd", "Mcr", "Atn", "Edn" ]
            return self.retrieve_token_and_filter(skips, "Txbx", *args)

        def header_textbox(self, *args):
            skips = [ "Text", "Ftn", "Hdd", "Mcr", "Atn", "Edn", "Txbx" ]
            return self.retrieve_token_and_filter(skips, "HdrTxbx", *args)

        def parse_fib(self, f):
            if self.get_ushort(f, 0x0000) != 0xa5ec:
                raise ValueError('Not a Word document')

            nFib = self.get_ushort(f, 0x0002)
            if nFib < 101:
                raise ValueError('Unsupported version')

            flags = self.get_ushort(f, 0x000a)

            self.flag["fComplex"] = (flags & 0x0004 != 0)

            self.flag["fEncrypted"] = (flags & 0x0100 != 0)
            if self.flag["fEncrypted"]:
                raise ValueError('Encrypted MSWord document file is not supported')

            self.flag["fWhichTblStm"] = (flags & 0x0200 != 0)

            self.fcMin = self.get_ulong(f, 0x0018)
            self.fcMac = self.get_ulong(f, 0x001c)
            self.cbMac = self.get_ulong(f, 0x0040)

            self.fcClx  = self.get_ulong(f, 0x01a2)
            self.lcbClx = self.get_ulong(f, 0x01a6)

            self.parse_fib_ccps(f)

        def parse_fib_ccps(self, f):
            for key, offset in __class__.OFFSET_FIB_CCP_MAP.items():
                self.ccp[key] = self.get_ulong(f, offset)

        def parse_piece_table(self, f):
            if self.lcbClx <= 0:
                ccpAll = sum(self.ccp[key] for key in __class__.OFFSET_FIB_CCP_MAP)

                self.pcds = [{
                    "fc": self.fcMin,
                    "cp": 0,
                    "ccp": ccpAll,
                }]
                return

            f.seek(self.fcClx)
            clx = BytesIO(f.read(self.lcbClx))

            while clx.tell() < self.lcbClx:
                clxt = struct.unpack("B", clx.read(1))[0]
                if clxt == 2:
                    break

                if clxt == 1:
                    skip = struct.unpack('H', clx.read(2))[0]
                    clx.seek(skip)
                else:
                    raise ValueError('Unknown CLX block')

            if clx.tell() == self.lcbClx:
                raise ValueError('PCDs not found')

            length = struct.unpack('I', clx.read(4))[0]

            n = int((length - 4) / (4 + 8))
            cps = [struct.unpack('I', clx.read(4))[0] for _ in range(n+1)]

            self.pcds = []
            for i in range(1, n+1):
                pcd_data = clx.read(8)
                fc = struct.unpack('H', pcd_data[2:4])[0]

                self.pcds.append({
                "fc" : fc,
                "cp"  : cps[i - 1],
                "ccp" : cps[i] - cps[i - 1],
                })

        def retrieve_substring(self, f, offset, length = -1):
            i = 0
            while i < len(self.pcds):
                if self.pcds[i]["cp"] > offset:
                    break
                i += 1
            i -= 1
            if i < 0:
                raise ValueError('could not find suitable heading piece')

            output = b""

            while length > 0 or length < 0:
                pcd = self.pcds[i]

                _len = length
                if pcd["ccp"] < _len or _len < 0:
                    _len = pcd["ccp"]

                pcdfc = pcd["fc"]
                if pcdfc == 0x40000000:
                    # cp1252
                    fc = (pcdfc ^ 0x40000000) >> 1
                    fc += offset
                    offset = 0

                    f.seek(fc)
                    output += f.read(_len).decode("cp1252").encode("utf8")
                else:
                    # utf-16-le
                    fc = pcd["fc"]
                    fc += offset * 2
                    offset = 0
                    f.seek(fc)
                    output += f.read(_len * 2).decode("utf-16-le").encode("utf8")

                if length >= 0:
                    length -= _len

                i += 1
                if i >= len(self.pcds):
                    break

            return output

        def get_ushort(self, f, pos):
            f.seek(pos)
            return struct.unpack("H", f.read(2))[0]

        def get_ulong(self, f, pos):
            f.seek(pos)
            return struct.unpack("I", f.read(4))[0]

        def retrieve_token_and_filter(self, skip_tokens, target, *args):
            skip = 0

            for key in skip_tokens:
                skip += self.ccp[key]

            return self.retrieve_and_filter(skip, self.ccp[target], *args)

        def retrieve_and_filter(self, offset, length, *args):

            string = self.retrieve_substring(self.ole.openstream('WordDocument'), offset, length)

            if "raw" not in args:
                return  self.format_into_plain(string)

            return string

        def format_into_plain(self, dat):
            dat = __class__.re_formatbin.sub(b'\1\n', dat)
            reps = [
                  (b"\x0d" , b"\n"),         # ASIS: Line Feed
                  (b"\x09" , b"\t"),         # ASIS: Tab
                  (b"\x0b" , b"\n"),         # Hard line breaks
                  (b"\x2d" , b"\x2d"),       # ASIS, breaking hyphens; U+2010?
                  (b"\x1f" , b"\xad"),   # Non-required hyphens (into Soft hyphen)
                  (b"\x1e" , b"\x11 "),   # Non-breaking hyphens
                  (b"\xa0" , b"\xa0"),       # ASIS: Non-breaking-spaces
                  (b"\x0c" , b"\x0c"),       # ASIS: Page breaks or Section marks
                  (b"\x0e" , b"\x0e"),       # ASIS: Column breaks
                  (b"\x13" , b""),           # Field begin mark
                  (b"\x15" , b""),           # Field end mark
                  (b"\x14" , b""),           # Field separator
                  (b"\x07" , b"\t"),         # Cell mark or Row mark
                  (b"\x01", b""),            # table record delimiter
            ]
            for k, v in reps:
                dat = dat.replace(k, v)
            return dat
except ModuleNotFoundError:
    olefile = MSWordOLE = NotInstalledModuleError("Please Install command: pip3 install olefile")

try:
    from pdfminer.pdfparser import PDFParser, PDFDocument
    from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
    from pdfminer.converter import PDFPageAggregator
    from pdfminer.layout import LAParams, LTTextBox, LTTextLine
except ModuleNotFoundError:
    PDFParser = PDFDocument = PDFResourceManager = PDFPageInterpreter = PDFPageAggregator = \
    LAParams = LTTextBox = LTTextLine = NotInstalledModuleError("Please Install command: pip3 install pdfminer3k")

try:
    from pyperclip import paste as getclip, copy as setclip
except ModuleNotFoundError:
    if os.name == "nt":
        import tkinter as tk
        def getclip():
            a=tk.Tk()
            return a.clipboard_get()

        import subprocess
        def setclip(text):
            if not isinstance(text, (str, int, float, bool)):
                raise RuntimeError('only str, int, float, and bool values can be copied to the clipboard, not %s' % (text.__class__.__name__))
            text = str(text)
            p = subprocess.Popen(['clip.exe'],
                                 stdin=subprocess.PIPE, close_fds=True)
            p.communicate(input=text.encode("cp932"))

    else:
        getclip = NotInstalledModuleError("Please Install command: pip3 install pyperclip")
        setclip = NotInstalledModuleError("Please Install command: pip3 install pyperclip")

try:
    import pylnk
    pylnkfile = pylnk.file
except ModuleNotFoundError:
    pylnkfile = NotInstalledModuleError("Please Install command pip3 install pylnk")

try:
    from sqlalchemy import create_engine, Table, MetaData
except:
    create_engine = Table = MetaData = NotInstalledModuleError("Please Install command: pip3 install sqlalchemy")

try:
    import pyodbc
except:
    if os.name == 'posix':
        pyodbc = NotInstalledModuleError("Please Install command: \nubuntu => apt-get install unixodbc-dev;pip3 install pyodbc\nredhat => yum install unixODBC unixODBC-devel;pip3 install pyodbc\n")
    else:
        pyodbc = NotInstalledModuleError("Please Install command: pip3 install pyodbc")

from util.filetype import guesstype
from util.core import binopen, opener, getencoding, binchunk, globbing

def lsdir(path, recursive=True):
    func = "rglob" if recursive else "glob"
    for p in map(pathlib.Path, glob(str(path))):
        yield p
        for r in p.__getattribute__(func)("*"):
            yield r

def pathbin(path_or_buffer):
    fp = binopen(path_or_buffer)
    try:
        path = pathlib.Path(fp.name)
        return path, fp
    except AttributeError:
        return fp._fp.name, fp
    except TypeError:
        return None, fp

class _HTMLParser(HTMLParser):
    def __init__(self, data, encoding=None):
        super().__init__()
        self.result = []
        self.feed(data, encoding)

    def handle_data(self, data):
        self.result.append(data)

    def feed(self, data, encoding=None):
        if isinstance(data, bytes):
            if encoding:
                data = data.decode(encoding)
            else:
                data = data.decode()
        super().feed(data)

class _XMLHandler(object):
    def __init__(self,
                 item_depth=0,
                 item_callback=lambda *args: True,
                 xml_attribs=True,
                 attr_prefix='@',
                 cdata_key='#text',
                 force_cdata=False,
                 cdata_separator='',
                 postprocessor=None,
                 dict_constructor=dict,
                 strip_whitespace=True,
                 nssep=':',
                 namespaces=None,
                 force_list=None,
                 comment_key='#comment'):
        self.path = []
        self.stack = []
        self.data = []
        self.item = None
        self.item_depth = item_depth
        self.xml_attribs = xml_attribs
        self.item_callback = item_callback
        self.attr_prefix = attr_prefix
        self.cdata_key = cdata_key
        self.force_cdata = force_cdata
        self.cdata_separator = cdata_separator
        self.postprocessor = postprocessor
        self.dict_constructor = dict_constructor
        self.strip_whitespace = strip_whitespace
        self.nssep = nssep
        self.namespaces = namespaces
        self.namespace_declarations = dict()
        self.force_list = force_list
        self.comment_key = comment_key

    def _build_name(self, full_name):
        if self.namespaces is None:
            return full_name
        i = full_name.rfind(self.nssep)
        if i == -1:
            return full_name
        namespace, name = full_name[:i], full_name[i+1:]
        try:
            short_namespace = self.namespaces[namespace]
        except KeyError:
            short_namespace = namespace
        if not short_namespace:
            return name
        else:
            return self.nssep.join((short_namespace, name))

    def _attrs_to_dict(self, attrs):
        if isinstance(attrs, dict):
            return attrs
        return self.dict_constructor(zip(attrs[0::2], attrs[1::2]))

    def startNamespaceDecl(self, prefix, uri):
        self.namespace_declarations[prefix or ''] = uri

    def startElement(self, full_name, attrs):
        name = self._build_name(full_name)
        attrs = self._attrs_to_dict(attrs)
        if attrs and self.namespace_declarations:
            attrs['xmlns'] = self.namespace_declarations
            self.namespace_declarations = dict()
        self.path.append((name, attrs or None))
        if len(self.path) > self.item_depth:
            self.stack.append((self.item, self.data))
            if self.xml_attribs:
                attr_entries = []
                for key, value in attrs.items():
                    key = self.attr_prefix+self._build_name(key)
                    if self.postprocessor:
                        entry = self.postprocessor(self.path, key, value)
                    else:
                        entry = (key, value)
                    if entry:
                        attr_entries.append(entry)
                attrs = self.dict_constructor(attr_entries)
            else:
                attrs = None
            self.item = attrs or None
            self.data = []

    def endElement(self, full_name):
        name = self._build_name(full_name)
        if len(self.path) == self.item_depth:
            item = self.item
            if item is None:
                item = (None if not self.data
                        else self.cdata_separator.join(self.data))

            should_continue = self.item_callback(self.path, item)
            if not should_continue:
                raise InterruptedError
        if len(self.stack):
            data = (None if not self.data
                    else self.cdata_separator.join(self.data))
            item = self.item
            self.item, self.data = self.stack.pop()
            if self.strip_whitespace and data:
                data = data.strip() or None
            if data and self.force_cdata and item is None:
                item = self.dict_constructor()
            if item is not None:
                if data:
                    self.push_data(item, self.cdata_key, data)
                self.item = self.push_data(self.item, name, item)
            else:
                self.item = self.push_data(self.item, name, data)
        else:
            self.item = None
            self.data = []
        self.path.pop()

    def characters(self, data):
        if not self.data:
            self.data = [data]
        else:
            self.data.append(data)

    def comments(self, data):
        if self.strip_whitespace:
            data = data.strip()
        self.item = self.push_data(self.item, self.comment_key, data)

    def push_data(self, item, key, data):
        if self.postprocessor is not None:
            result = self.postprocessor(self.path, key, data)
            if result is None:
                return item
            key, data = result
        if item is None:
            item = self.dict_constructor()
        try:
            value = item[key]
            if isinstance(value, list):
                value.append(data)
            else:
                item[key] = [value, data]
        except KeyError:
            if self._should_force_list(key, data):
                item[key] = [data]
            else:
                item[key] = data
        return item

    def _should_force_list(self, key, value):
        if not self.force_list:
            return False
        if isinstance(self.force_list, bool):
            return self.force_list
        try:
            return key in self.force_list
        except TypeError:
            return self.force_list(self.path[:-1], key, value)


def xmltodict(data:bytes, encoding=None, namespaces=False, nssep=':', comments=False, **kw):
    hl = _XMLHandler(nssep=nssep, **kw)
    ps = ParserCreate(encoding, nssep or None)
    ps.ordered_attributes = True
    ps.StartNamespaceDeclHandler = hl.startNamespaceDecl
    ps.StartElementHandler = hl.startElement
    ps.EndElementHandler = hl.endElement
    ps.CharacterDataHandler = hl.characters

    if comments:
        ps.CommentHandler = hl.comments

    ps.buffer_text = True
    ps.Parse(data, True)

    return hl.item

def _iterdecoder(fp, encoding=None, buffer=1024**2):
    if encoding:
        while True:
            buf = fp.read(buffer) + fp.readline()
            if not buf:
                break
            for line in buf.decode(encoding).splitlines():
                yield line
    else:
        while True:
            buf = fp.read(buffer) + fp.readline()
            if not buf:
                break
            for line in buf.decode().splitlines():
                yield line

pinfo = namedtuple("LazyReader", ["path", "target", "value"])
class readrow:
    @classmethod
    def handler(cls, path_or_buffer):
        fp = binopen(path_or_buffer)
        funcstr = guesstype(fp)

        if funcstr is None:
            raise NotImplementedError("Unknown binary data `{}`".format(fp))
        try:
            return getattr(cls, funcstr)
        except AttributeError:
            raise NotImplementedError("{} class is not function `{}`".format(cls.__name__, funcstr))

    def __new__(cls, path_or_buffer, *args, **kw):
        return cls.handler(path_or_buffer)(path_or_buffer, *args, **kw)

    @staticmethod
    def pptx(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        for i, s in enumerate(Presentation(fp).slides):
            for t in (r.text for sp in s.shapes if sp.has_text_frame for p in sp.text_frame.paragraphs for r in p.runs if r.text):
                yield pinfo(path, i, t)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def ppt(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        mgl.textdump = b""
        mgl.params.noStructOutput = True
        mgl.params.dumpText = True
        mgl.params.noRawDumps = True
        mgl.params.showSectorChain = True

        strm = PPTFile(fp.read(), mgl.params)
        dirstrm = strm.getDirectoryStreamByName(b"PowerPoint Document")
        dirstrm.readRecords()

        for x in mgl.textdump.split(b"*\n")[-1].replace(b"\r", b"\n").decode().split("\n"):
            yield pinfo(path, None, x) #TODO None

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def docx(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        i = 0
        it = Document(fp).element.getiterator()
        for txt in it:
            if isinstance(txt, CT_Br):
                i += 1
            if isinstance(txt, CT_Text):
                yield pinfo(path, i, txt.text)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def doc(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        word = MSWordOLE(fp)
        for line in word.document().decode().splitlines():
            yield pinfo(path, None, line) #TODO None

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def xlsx(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        with xlrd.open_workbook(file_contents=fp.read()) as wb:
            for sh in wb.sheets():
                sname = sh.name
                for i in range(sh.nrows):
                    yield pinfo(path, sname, sh.row_values(i))

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    xls = xlsx

    @staticmethod
    def pdf(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        ps = PDFParser(fp)
        doc = PDFDocument()
        ps.set_document(doc)
        doc.set_parser(ps)
        doc.initialize('')

        mgr = PDFResourceManager(caching=True)
        dev = PDFPageAggregator(mgr, laparams=LAParams())
        ip = PDFPageInterpreter(mgr, dev)

        for i, page in enumerate(doc.get_pages(), 1):
            ip.process_page(page)
            text = "".join(x.get_text() for x in dev.get_result() if isinstance(x, (LTTextBox, LTTextLine)))
            yield pinfo(path, i, text.rstrip("{}\n".format(i)))

        dev.close()
        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def csv(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        dat = fp.read(SAMPLEBYTE) + fp.readline()

        e = getencoding(dat)
        txt = dat.decode(e)
        dialect = _csv.Sniffer().sniff(txt)

        for row in _csv.reader(StringIO(txt), dialect=dialect):
            yield pinfo(path, None, row)

        it = _iterdecoder(fp, e)
        for row in _csv.reader(it, dialect=dialect):
            yield pinfo(path, None, row)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def dml(path_or_buffer, excludes = ["^;?\n", "^;$", "^\s*//.*$"], lvsep = "  "):
        path, fp = pathbin(path_or_buffer)
        ret = []
        dat = fp.read()
        e = getencoding(dat)

        lines = list(enumerate((dat.decode(e) if e else dat.decode()).splitlines(),1))
        maxlen = max(x.count(lvsep) for i,x in lines)

        rec = [""] * maxlen

        for i, line in reversed(lines):

            if line and all(not re.match(x, line) for x in excludes):
                row = line.rstrip().split(lvsep)
                name = row[-1]
                if name.startswith("end"):
                    recordname = name.replace("end", "").strip("[; ]")
                    rec[row.index(name)] = recordname
                elif name.startswith("record"):
                    for x in range(row.index(name), maxlen):
                        rec[x] = ""
                else:
                    ret.append([i, ".".join(rec).strip("."), line.strip()])

        for i, rec, line in reversed(ret):
            yield pinfo(path, rec, line)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def lnk(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        r = pylnkfile()
        r.open_file_object(fp)
        for x in dir(r):
            attr = getattr(r, x)
            if x[0] == "_" or callable(attr):
                continue
            yield pinfo(path, x, attr)

    @staticmethod
    def txt(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        dat = fp.read(SAMPLEBYTE)
        fp.close()
        encoding = getencoding(dat)
        with codecs.open(path, encoding=encoding) as f:
            for line in f:
                yield pinfo(path, None, line.rstrip())
        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def zip(path_or_buffer, targets=[]):
        path, fp = pathbin(path_or_buffer)

        with ZipFile(fp) as f:
            if not targets:
                targets = f.namelist()
            for info in f.infolist():
                if info.is_dir():
                    continue
                target = info.filename.encode("cp437").decode("cp932")
                if target in targets:
                    fp = f.open(info)
                    func = __class__.handler(fp)
                    fp.seek(0)
                    for row in func(fp):
                        yield pinfo(path, target, row.value)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def gz(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        target = path.stem
        with GzipFile(fileobj=fp) as f:
            func = __class__.handler(f)
            f.seek(0)
            for row in func(f):
                yield pinfo(path, target, row)
        if not hasattr(path_or_buffer, "close"):
            fp.close()


    @staticmethod
    def bz2(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        target = path.stem
        with BZ2File(fp) as f:
            func = __class__.handler(f)
            f.seek(0)
            for row in func(f):
                yield pinfo(path, target, row)
        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def xz(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        target = path.stem
        with LZMAFile(fp) as f:
            func = __class__.handler(f)
            f.seek(0)
            for row in func(f):
                yield pinfo(path, target, row)
        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def tar(path_or_buffer, mode="r", targets=[]):
        path, fp = pathbin(path_or_buffer)

        with TarFile.open(mode=mode, fileobj=fp) as f:
            if not targets:
                targets = f.getnames()
            for info in f.getmembers():
                if info.isdir():
                    continue
                target = info.name

                if target in targets:
                    fp = f.extractfile(info)
                    func = __class__.handler(fp)
                    fp.seek(0)
                    for row in func(fp):
                        yield pinfo(path, target, row.value)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def lha(path_or_buffer, targets=[]):
        path, fp = pathbin(path_or_buffer)

        f = LhaFile(fp)
        if not targets:
            targets = f.namelist()
        for info in f.infolist():
            target = info.filename
            if target in targets:

                dat = f.read(info.filename)
                func = __class__.handler(dat)
                bio = BytesIO(dat)
                bio.name = target
                for row in func(bio):
                    yield pinfo(path, target, row.value)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def rar(path_or_buffer, targets=[]):
        path, fp = pathbin(path_or_buffer)

        with RarFile(fp) as f:
            if not targets:
                targets = f.namelist()
            for info in f.infolist():
                if info.isdir():
                    continue
                target = info.filename
                if target in targets:
                    fp = f.open(info)
                    func = __class__.handler(fp)
                    fp.seek(0)
                    for row in func(fp):
                        yield pinfo(path, target, row.value)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def locate(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        fp.read(8) # magic number
        conf_size = struct.unpack(">i", fp.read(4))[0]
        fp.read(4 + conf_size) # dbconf

        for entry in binchunk(fp, sep=b"\x00\x02"):
            raw = entry[16:].rstrip()
            raw = raw.replace(b"\x00\x00", b"\vfile\t").replace(b"\x00\x01", b"\vdir\t")

            root = ""
            for e in raw.decode("utf-8").split("\v"):
                try:
                    tp, name = e.split("\t")
                    if name == "/":
                        tp = "dir"
                    else:
                        name = root + "/" + name
                        if name in ("/etc/init.d", "/etc/rc0.d", "/etc/rc1.d", "/etc/rc2.d", "/etc/rc3.d", "/etc/rc4.d", "/etc/rc5.d", "/etc/rc6.d", "/etc/ssl/certs", "/etc/xdg/systemd/user", "/lib", "/lib64", "/sbin", "/usr/lib64/go/4.8.5", "/usr/libexec/gcc/x86_64-redhat-linux/4.8.5", "/bin", "/usr/include/c++/4.8.5", "/usr/lib/debug/bin", "/usr/lib/debug/lib", "/usr/lib/debug/lib64", "/usr/lib/debug/sbin", "/usr/lib/gcc/x86_64-redhat-linux/4.8.5", "/usr/lib/go/4.8.5", "/usr/lib/terminfo", "/usr/share/doc/git-1.8.3.1/contrib/hooks", "/usr/share/doc/redhat-release", "/usr/share/doc/vim-common-7.4.160/docs", "/usr/share/gcc-4.8.5", "/usr/share/gccxml-0.9/GCC/5.0", "/usr/share/gccxml-0.9/GCC/5.1", "/usr/share/gccxml-0.9/GCC/5.2", "/usr/share/gccxml-0.9/GCC/5.3", "/usr/share/gdb/auto-load/lib64", "/usr/share/groff/current", "/usr/tmp", "/var/lock", "/var/mail", "/var/run"):
                            tp = "dir"

                    yield pinfo(path, tp, name)

                except ValueError:
                    root = e

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def json(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        yield pinfo(path, None, jsonlib.loads(fp.read()))

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def html(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        dat = fp.read()
        if not hasattr(path_or_buffer, "close"):
            fp.close()
        p = _HTMLParser(dat, getencoding(dat))

        for r in p.result:
            yield pinfo(path, None, r)

        del p

    @staticmethod
    def xml(
        path_or_buffer_or_elem,
        parenttag=None,
        na_val=None,
        return_docinfo=True,
        comment=True,
        callback_attr = lambda x: ", ".join(["{}={}".format(k, repr(v)) for k, v in x]),
        dinfo = namedtuple("docinfo", ['doctype', 'encoding', 'externalDTD', 'internalDTD', 'public_id', 'root_name', 'standalone', 'system_url', 'xml_version']),
        einfo = namedtuple("elementinfo", ['tagpath', 'attr', 'value']),
        filename = None,
        ):

        if parenttag is None:
            parenttag = [""]

        if isinstance(path_or_buffer_or_elem, Element):
            e = path_or_buffer_or_elem
        else:
            try:
                filename, fp = pathbin(path_or_buffer_or_elem)
                tree = etree.parse(fp)
            except FileNotFoundError:
                tree = etree.fromstring(path_or_buffer_or_elem)

            if return_docinfo:
                docinfo = tree.docinfo

                yield pinfo(filename, "/", dinfo(*[getattr(docinfo, x) for x in dinfo._fields]))

                if comment:
                    for c in tree.xpath("/comment()"):
                        yield pinfo(filename, "/#comment", einfo("/", None, c.text))

            e = tree.getroot()

        parenttag.append("#comment" if e.tag is etree.Comment else e.tag)
        pt = "/".join(parenttag)

        try:
            val = e.text.replace("\xa0", "")
        except AttributeError:
            val = na_val

        attr = e.items() or na_val
        if attr and callback_attr:
            attr = callback_attr(attr)

        yield pinfo(filename, pt, einfo(pt, attr, val))  #, (e.sourceline, etree.tostring(e)))

        for child in e:
            if not comment and child.tag is etree.Comment:
                continue

            for r in __class__.xml(child, parenttag.copy(), filename=filename):
                yield r

        if not isinstance(path_or_buffer_or_elem, Element) and not hasattr(path_or_buffer_or_elem, "close"):
            fp.close()


    @staticmethod
    def mdb(path_or_buffer, targets=[], uid="", passwd=""):
        path, fp = pathbin(path_or_buffer)
        fp.close()
        try:
            driver = next(x for x in pyodbc.drivers() if x.startswith("Microsoft Access Driver "))
        except StopIteration:
            raise RuntimeError("Not Installed Microsoft Access Driver")

        dsnstr = r'DRIVER={{{}}};DBQ={};UID="{}";PWD="{}";'
        with pyodbc.connect(dsnstr.format(driver, path.resolve(), uid, passwd)) as con:
            with con.cursor() as cur:
                if not targets:
                    targets = [t.table_name for t in cur.tables() if t.table_name[:4] not in ["MSys", "~TMP"]]

            for table in targets:
                #columns = Table(table, MetaData(), autoload=True, autoload_with=engine).columns
                for row in con.execute("SELECT * FROM "+ table):
                    yield pinfo(path, table, list(row))

    accdb = mdb

    @staticmethod
    def sqlite3(path_or_buffer, targets=[]):
        path, fp = pathbin(path_or_buffer)
        fp.close()
        engine = create_engine(r"sqlite:///{}".format(path))
        path = engine.url
        if not targets:
            targets = engine.table_names()

        with engine.connect() as con:
            for table in targets:
                #columns = Table(table, MetaData(), autoload=True, autoload_with=engine).columns
                for row in con.execute("SELECT * FROM "+ table):
                    yield pinfo(path, table, list(row))

    @staticmethod
    def pickle(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        yield pinfo(path, None, pkload(fp))

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def clipboard(ftype="csv"):
        """
        clipboard read parse

        ftype : xsv or html, xml, txt, json
        """
        funcs = getattr(__class__, ftype)

        cdat = getclip()
        with BytesIO() as bio:
            try:
                bdat = cdat.encode("utf-8")
            except:
                bdat = cdat.encode()
            bio.write(bdat)
            bio.seek(0)
            bio.name = "clipboad"
            for r in funcs(bio):
                yield r

dumper = globbing(readrow)

class grouprow(readrow):

    @staticmethod
    def pptx(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        for i, s in enumerate(Presentation(fp).slides):
            yield pinfo(path, i, [sp.text for sp in s.shapes if sp.has_text_frame])

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def ppt(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        mgl.textdump = b""
        mgl.params.noStructOutput = True
        mgl.params.dumpText = True
        mgl.params.noRawDumps = True
        mgl.params.showSectorChain = True

        strm = PPTFile(fp.read(), mgl.params)
        dirstrm = strm.getDirectoryStreamByName(b"PowerPoint Document")
        dirstrm.readRecords()

        yield pinfo(path, None, mgl.textdump.split(b"*\n")[-1].replace(b"\r", b"\n").decode().split("\n")) #TODO None

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def docx(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        i = 0
        it = Document(fp).element.getiterator()
        ret = []
        for txt in it:
            if isinstance(txt, (CT_Br, CT_PageMar)):
                yield pinfo(path, i, ret)
                ret = []
                i += 1
            if isinstance(txt, CT_Text):
                ret.append(txt.text)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def doc(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        word = MSWordOLE(fp)
        yield pinfo(path, None, word.document().decode().splitlines()) #TODO None

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def xlsx(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        with xlrd.open_workbook(file_contents=fp.read()) as wb:
            for sh in wb.sheets():
                yield pinfo(path, sh.name, [sh.row_values(i) for i in range(sh.nrows)])

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    xls = xlsx

    @staticmethod
    def pdf(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        ps = PDFParser(fp)
        doc = PDFDocument()
        ps.set_document(doc)
        doc.set_parser(ps)
        doc.initialize('')

        mgr = PDFResourceManager(caching=True)
        dev = PDFPageAggregator(mgr, laparams=LAParams())
        ip = PDFPageInterpreter(mgr, dev)

        for i, page in enumerate(doc.get_pages(), 1):
            ip.process_page(page)
            rows = [x.get_text() for x in dev.get_result() if isinstance(x, (LTTextBox, LTTextLine))]
            yield pinfo(path, i, rows)

        dev.close()
        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def csv(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        dat = fp.read(SAMPLEBYTE) + fp.readline()

        e = getencoding(dat)
        txt = dat.decode(e)
        dialect = _csv.Sniffer().sniff(txt)

        tmp = list(_csv.reader(StringIO(txt), dialect=dialect))

        it = _iterdecoder(fp, e)
        yield pinfo(path, None, tmp + list(_csv.reader(it, dialect=dialect)))

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def txt(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        dat = fp.read(SAMPLEBYTE)
        fp.close()
        encoding = getencoding(dat)
        with codecs.open(path, encoding=encoding) as f:
            yield pinfo(path, None, f.readlines())

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def dml(path_or_buffer, excludes = ["^;?\n", "^;$", "^\s*//.*$"], lvsep = "  "):
        path, fp = pathbin(path_or_buffer)
        ret = []
        dat = fp.read()
        e = getencoding(dat)

        lines = list(enumerate((dat.decode(e) if e else dat.decode()).splitlines(),1))
        maxlen = max(x.count(lvsep) for i,x in lines)

        rec = [""] * maxlen

        for i, line in reversed(lines):

            if line and all(not re.match(x, line) for x in excludes):
                row = line.rstrip().split(lvsep)
                name = row[-1]
                if name.startswith("end"):
                    recordname = name.replace("end", "").strip("[; ]")
                    rec[row.index(name)] = recordname
                elif name.startswith("record"):
                    for x in range(row.index(name), maxlen):
                        rec[x] = ""
                else:
                    ret.append([i, ".".join(rec).strip("."), line.strip()])

        gr = {}
        for i, rec, line in reversed(ret):
            if rec in gr:
                gr[rec].append(line)
            else:
                gr[rec] = [line]

        for rec, lines in gr.items():
            yield pinfo(path, rec, lines)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def lnk(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        r = pylnkfile()
        r.open_file_object(fp)
        tar = r.environment_variables_location
        ret = {}
        for x in dir(r):
            attr = getattr(r, x)
            if x[0] == "_" or callable(attr):
                continue
            ret[x] = attr
        yield pinfo(path, tar, ret)

    @staticmethod
    def html(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        dat = fp.read()
        if not hasattr(path_or_buffer, "close"):
            fp.close()
        p = _HTMLParser(dat, getencoding(dat))
        yield pinfo(path, None, p.result)

        del p

    @staticmethod
    def mdb(path_or_buffer, targets=[], uid="", passwd=""):
        path, fp = pathbin(path_or_buffer)
        fp.close()
        try:
            driver = next(x for x in pyodbc.drivers() if x.startswith("Microsoft Access Driver "))
        except StopIteration:
            raise RuntimeError("Not Installed Microsoft Access Driver")

        dsnstr = r'DRIVER={{{}}};DBQ={};UID="{}";PWD="{}";'
        with pyodbc.connect(dsnstr.format(driver, path.resolve(), uid, passwd)) as con:
            with con.cursor() as cur:
                if not targets:
                    targets = [t.table_name for t in cur.tables() if t.table_name[:4] not in ["MSys", "~TMP"]]

            for table in targets:
                #columns = Table(table, MetaData(), autoload=True, autoload_with=engine).columns
                yield pinfo(path, table, list(map(list, con.execute("SELECT * FROM "+ table))))

    accdb = mdb

    @staticmethod
    def sqlite3(path_or_buffer, targets=[]):
        path, fp = pathbin(path_or_buffer)
        fp.close()
        engine = create_engine(r"sqlite:///{}".format(path))
        path = engine.url
        if not targets:
            targets = engine.table_names()

        with engine.connect() as con:
            for table in targets:
                #columns = Table(table, MetaData(), autoload=True, autoload_with=engine).columns
                yield pinfo(path, table, list(map(list, con.execute("SELECT * FROM "+ table))))

    @staticmethod
    def xml(path_or_buffer):
        ret = {}
        for path, target, value in readrow.xml(path_or_buffer):
            if target in ret:
                ret[target].append(value)
            else:
                ret[target] = [value]

        for k, v in ret.items():
            yield pinfo(path, k, v)

@lru_cache()
def ts2date(x):#, dfm = "%Y/%m/%d %H:%M"):
    return datetime.fromtimestamp(x)#.strftime(dfm)

@lru_cache()
def oct2perm(x):
    return oct(x)[-3:]

if os.name == "posix":
    from pwd import getpwuid
    from grp import getgrgid

    @lru_cache()
    def getuser(uid):
        try:
            return getpwuid(uid).pw_name
        except KeyError:
            return uid

    @lru_cache()
    def getgroup(gid):
        try:
            return getgrgid(gid).gr_name
        except KeyError:
            return gid
else:
    def getuser(uid): return uid
    def getgroup(gid): return gid

def getsize_by_seek(fp, chunksize=1024**2):
    try:
        fp.seek(0, 2)
    except OSError:
        pos = -1
        while fp.tell() != pos:
            pos = fp.tell()
            fp.seek(pos + chunksize)

    return fp.tell()

def getsize_gz(fp):
    gz = fp.fileobj
    orggz = gz.tell()
    p = os.stat(gz.name).st_size

    if p < 402653184:
        gz.seek(-4, 2)
        c = sum(gz.read(4))
        gz.seek(orggz)
        if p < c:
            return c

    return getsize_by_seek(fp)


def getLF(dat:bytes, lf=re.compile(b"\r?\n|\r")):
    r = lf.search(dat)
    if r:
        return r.group(0).decode()
    return


sinfo = namedtuple("LazyInfo", ["path", "target", "parentdir", "basename", "is_file", "ext", "encoding", "LF", "owner", "group", "permission", "mtime", "ctime", "atime", "size"])
class getinfo:

    @classmethod
    def handler(cls, path_or_buffer):
        fp = binopen(path_or_buffer)
        func = guesstype(fp)
        if func is None:
            return cls.binaryfile
        elif func in [fc for fc in dir(cls) if not fc.startswith("_")]:
            return getattr(cls, func)
        else:
            return cls.flatfile

    def __new__(cls, path_or_buffer, *args, **kw):
        return cls.handler(path_or_buffer)(path_or_buffer, *args, **kw)

    @staticmethod
    def flatfile(path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        dat = fp.read(SAMPLEBYTE) + fp.readline()

        if not hasattr(path_or_buffer, "close"):
            fp.close()

        e = getencoding(dat)
        lf = getLF(dat)
        if hasattr(fp, "stat"):
            stat = fp.stat.copy()
            stat[6] = e
            stat[7] = lf
            return sinfo(*stat)
        else:
            stat = path.stat()
            return sinfo(
                path,
                None,
                path.parent,
                path.stem,
                path.is_file(),
                path.suffix,
                e,
                lf,
                getuser(stat.st_uid),
                getgroup(stat.st_gid),
                oct2perm(stat.st_mode),
                ts2date(stat.st_mtime),
                ts2date(stat.st_ctime),
                ts2date(stat.st_atime),
                stat.st_size
            )

    @staticmethod
    def binaryfile(path_or_buffer):
        path, fp = pathbin(path_or_buffer)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

        if hasattr(fp, "stat"):
            return sinfo(*fp.stat)
        else:
            stat = path.stat()
            return sinfo(
                path,
                None,
                path.parent,
                path.stem,
                path.is_file(),
                path.suffix,
                None,
                None,
                getuser(stat.st_uid),
                getgroup(stat.st_gid),
                oct2perm(stat.st_mode),
                ts2date(stat.st_mtime),
                ts2date(stat.st_ctime),
                ts2date(stat.st_atime),
                stat.st_size
            )

    @staticmethod
    def zip(path_or_buffer, targets=[]):
        yield __class__.binaryfile(path_or_buffer)

        path, fp = pathbin(path_or_buffer)

        with ZipFile(fp) as f:
            if not targets:
                targets = f.namelist()
            for info in f.infolist():
                target = info.filename.encode("cp437").decode("cp932")
                if target in targets:
                    zp = f.open(info)
                    func = __class__.handler(zp)
                    zp.seek(0)
                    zp.stat = [
                        path,  #path
                        target, #target
                        path,   #parentdir
                        target,    #basename
                        not info.is_dir(),    #is_file
                        os.path.splitext(target)[1],   #ext
                        None,   #encoding
                        None,   #LF
                        None,   #owner
                        None,   #group
                        None,   #permission
                        datetime(*info.date_time),    #mtime
                        None,   #ctime
                        None,   #atime
                        info.file_size   #size
                    ]
                    yield func(zp)


        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def tar(path_or_buffer, mode="r", targets=[]):
        yield __class__.binaryfile(path_or_buffer)

        path, fp = pathbin(path_or_buffer)

        with TarFile.open(mode=mode, fileobj=fp) as f:
            if not targets:
                targets = f.getnames()
            for info in f.getmembers():
                target = info.name

                if target in targets:
                    zp = f.extractfile(info)
                    func = __class__.handler(zp)
                    zp.seek(0)
                    zp.stat = [
                        path,  #path
                        target, #target
                        path,   #parentdir
                        target,    #basename
                        info.isfile(),    #is_file
                        os.path.splitext(target)[1],   #ext
                        None,   #encoding
                        None,   #LF
                        info.uname,   #owner
                        info.gname,   #group
                        oct2perm(info.mode),   #permission
                        ts2date(info.mtime),    #mtime
                        None,   #ctime
                        None,   #atime
                        info.size   #size
                    ]
                    yield func(zp)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def rar(path_or_buffer, targets=[]):
        yield __class__.binaryfile(path_or_buffer)

        path, fp = pathbin(path_or_buffer)

        with RarFile(fp) as f:
            if not targets:
                targets = f.namelist()
            for info in f.infolist():
                target = info.filename
                if target in targets:
                    zp = f.open(info)
                    func = __class__.handler(zp)
                    zp.seek(0)
                    zp.stat = [
                        path,  #path
                        target, #target
                        path,   #parentdir
                        target,    #basename
                        not info.isdir(),    #is_file
                        os.path.splitext(target)[1],   #ext
                        None,   #encoding
                        None,   #LF
                        None,   #owner
                        None,   #group
                        None,   #permission
                        datetime(*info.date_time),    #mtime
                        None,   #ctime
                        None,   #atime
                        info.file_size   #size
                    ]
                    yield func(zp)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def lha(path_or_buffer, targets=[]):
        yield __class__.binaryfile(path_or_buffer)

        path, fp = pathbin(path_or_buffer)

        f = LhaFile(fp)
        if not targets:
            targets = f.namelist()
        for info in f.infolist():
            target = info.filename
            if target in targets:
                zp = BytesIO(f.read(target))
                zp.name = info.filename
                func = __class__.handler(zp)
                zp.seek(0)
                zp.stat = [
                    path,  #path
                    target, #target
                    path,   #parentdir
                    target,    #basename
                    True,    #is_file
                    os.path.splitext(target)[1],   #ext
                    None,   #encoding
                    None,   #LF
                    None,   #owner
                    None,   #group
                    None,   #permission
                    info.date_time,    #mtime
                    None,   #ctime
                    None,   #atime
                    info.file_size   #size
                ]
                yield func(zp)
        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def gz(path_or_buffer):
        yield __class__.binaryfile(path_or_buffer)

        path, fp = pathbin(path_or_buffer)
        target = path.stem
        with GzipFile(fileobj=fp) as f:
            func = __class__.handler(f)
            f.seek(0)
            target = os.path.splitext(fp.name)[0]
            pardir, basename = os.path.split(target)

            f.stat = [
                        path,  #path
                        target, #target
                        pardir,   #parentdir
                        basename,    #basename
                        True,    #is_file
                        os.path.splitext(target)[1],   #ext
                        None,   #encoding
                        None,   #LF
                        None,   #owner
                        None,   #group
                        None,   #permission
                        ts2date(f.mtime),    #mtime
                        None,   #ctime
                        None,   #atime
                        getsize_gz(f)   #size
                    ]
            yield func(f)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def bz2(path_or_buffer):
        yield __class__.binaryfile(path_or_buffer)

        path, fp = pathbin(path_or_buffer)
        target = path.stem
        with BZ2File(fp) as f:
            func = __class__.handler(f)
            f.seek(0)
            target = os.path.splitext(fp.name)[0]
            pardir, basename = os.path.split(target)

            f.stat = [
                        path,  #path
                        target, #target
                        pardir,   #parentdir
                        basename,    #basename
                        True,    #is_file
                        os.path.splitext(target)[1],   #ext
                        None,   #encoding
                        None,   #LF
                        None,   #owner
                        None,   #group
                        None,   #permission
                        None,    #mtime
                        None,   #ctime
                        None,   #atime
                        getsize_by_seek(f)   #size
                    ]
            f.seek(0)
            yield func(f)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

    @staticmethod
    def xz(path_or_buffer):
        yield __class__.binaryfile(path_or_buffer)

        path, fp = pathbin(path_or_buffer)
        target = path.stem
        with LZMAFile(fp) as f:
            func = __class__.handler(f)
            f.seek(0)
            target = os.path.splitext(fp.name)[0]
            pardir, basename = os.path.split(target)

            f.stat = [
                        path,  #path
                        target, #target
                        pardir,   #parentdir
                        basename,    #basename
                        True,    #is_file
                        os.path.splitext(target)[1],   #ext
                        None,   #encoding
                        None,   #LF
                        None,   #owner
                        None,   #group
                        None,   #permission
                        None,    #mtime
                        None,   #ctime
                        None,   #atime
                        getsize_by_seek(f)   #size
                    ]
            f.seek(0)
            yield func(f)

        if not hasattr(path_or_buffer, "close"):
            fp.close()

class getsize:
    @classmethod
    def handler(cls, path_or_buffer):
        path, fp = pathbin(path_or_buffer)
        func = guesstype(fp)
        if func in [fc for fc in dir(cls) if not fc.startswith("_")]:
            return getattr(cls, func)
        elif path:
            return path.stat().st_size
        else:
            return getsize_by_seek(fp)

    def __new__(cls, path_or_buffer, *args, **kw):
        path, fp = pathbin(path_or_buffer)
        func = guesstype(fp)
        if func in [fc for fc in dir(cls) if not fc.startswith("_")]:
            return getattr(cls, func)(fp, *args, **kw)
        elif path:
            return path.stat().st_size
        else:
            return getsize_by_seek(fp)

    @staticmethod
    def zip(path_or_buffer):
        fp = ZipFile(binopen(path_or_buffer))
        return sum(i.file_size for i in fp.infolist())

    @staticmethod
    def tar(path_or_buffer, mode="r"):
        fp = TarFile.open(mode=mode, fileobj=binopen(path_or_buffer))
        return sum(i.size for i in fp.getmembers())

    @staticmethod
    def rar(path_or_buffer, targets=[]):
        fp = RarFile(binopen(path_or_buffer))
        return sum(i.file_size for i in fp.infolist())

    @staticmethod
    def lha(path_or_buffer, targets=[]):
        fp = LhaFile(binopen(path_or_buffer))
        return sum(i.file_size for i in fp.infolist())

    @staticmethod
    def gz(path_or_buffer):
        fp = GzipFile(fileobj=binopen(path_or_buffer))
        return getsize_gz(fp)

    @staticmethod
    def bz2(path_or_buffer):
        return getsize_by_seek(BZ2File(binopen(path_or_buffer)))

    @staticmethod
    def xz(path_or_buffer):
        return getsize_by_seek(LZMAFile(binopen(path_or_buffer)))

class DBrow:
    @classmethod
    def handler(cls, path_or_buffer):
        fp = binopen(path_or_buffer)
        funcstr = guesstype(fp)

        if funcstr is None:
            raise NotImplementedError("Unknown binary data `{}`".format(fp))
        try:
            return getattr(cls, funcstr)
        except AttributeError:
            raise NotImplementedError("{} class is not function `{}`".format(cls.__name__, funcstr))

    def __new__(cls, path_or_buffer, *args, **kw):
        return cls.handler(path_or_buffer)(path_or_buffer, *args, **kw)

    @staticmethod
    def readsql(eng, dbname, server, targets):
        if not targets:
            targets = eng.table_names()

        with eng.connect() as con:
            for table in targets:
                #columns = Table(table, MetaData(), autoload=True, autoload_with=engine).columns
                schema_table = "{}/{}".format(dbname, table)
                for row in con.execute("SELECT * FROM "+ table):
                    yield pinfo(server, schema_table, list(row))

    @staticmethod
    def oracle(server, dbname, uid, passwd, port="1521", targets=[]):
        eng = create_engine("oracle+cx_oracle://{}:{}@{}/{}?port={}".format(uid, passwd,server,dbname,port))
        return __class__.readsql(eng, dbname, server, targets)

    @staticmethod
    def sqlserver(server, dbname, uid, passwd, port=None, targets=[]):
        target = "SQL Server Native Client"
        driver = sorted(x for x in pyodbc.drivers() if x.startswith(target))
        if not driver:
            raise RuntimeError("Not Found ODBC Driver " + target)
        dsnstr = r'DRIVER={{{}};SERVER={};DATABASE={};UID={};PWD={}'
        uri = 'mssql+pyodbc:///?odbc_connect=' + quote_plus(dsnstr.format(driver[-1],server,dbname,uid,passwd))
        eng = create_engine(uri)
        return __class__.readsql(eng, dbname, server, targets)

    @staticmethod
    def mysql(server, dbname, uid, passwd, port="3309", targets=[]):
        eng = create_engine(r'mysql://{}:{}@{}:{}/{}'.format(uid,passwd,server,port,dbname))
        return __class__.readsql(eng, dbname, server, targets)

    @staticmethod
    def postgres(server, dbname, uid, passwd, port="5432", targets=[]):
        eng = create_engine(r'postgresql://{}:{}@{}:{}/{}'.format(uid,passwd,server,port,dbname))
        return __class__.readsql(eng, dbname, server, targets)

    @staticmethod
    def db2(server, dbname, uid, passwd, port="50000", targets=[]):
        eng = create_engine(r'db2+ibm_db://{}:{}@{}:{}/{}'.format(uid,passwd,server,port,dbname))
        return __class__.readsql(eng, dbname, server, targets)

    @staticmethod
    def vertica(server, dbname, uid, passwd, port="5433", targets=[]):
        eng = create_engine(r'vertica+vertica_python://{}:{}@{}:{}/{}'.format(uid, passwd, server, port, dbname))
        return __class__.readsql(eng, dbname, server, targets)

class DBgrouprow(DBrow):
    @staticmethod
    def readsql(eng, dbname, server, targets):
        if not targets:
            targets = eng.table_names()

        with eng.connect() as con:
            for table in targets:
                #columns = Table(table, MetaData(), autoload=True, autoload_with=engine).columns
                schema_table = "{}/{}".format(dbname, table)
                rows = map(list, con.execute("SELECT * FROM "+ table))
                yield pinfo(server, schema_table, list(rows))

class Path(type(pathlib.Path())):

    __slots__ = (
        '_accessor',
        '_closed',
        '_encoding',
        '_ext',
        '_info',
        '_fullpath',
        '_file',
        '_type',
        '_rows',
    )

    def _init(self, *args, **kw):
        super()._init(*args, **kw)
        self._fullpath = None
        self._rows = None
        self._type = None
        self._file = None
        self._info = None

    @property
    def fullpath(self):
        if self._fullpath is None:
            self._fullpath = self.__str__() #self.joinpath(self.__str__(), self.content)
        return self._fullpath

    @property
    def info(self):
        if self._info is None:
            self._info = getinfo.binaryfile(self)
        return self._info

    def read_bytes(self, n=-1):
        """
        Open the file in bytes mode, read it, and close the file.
        """
        with open(self.__str__(), mode='rb') as f:
            return f.read(n)

    def read_text(self, n=-1, encoding=None, errors=None):
        """
        Open the file in text mode, read it, and close the file.
        """
        with open(self.__str__(), mode='r', encoding=encoding or self.encoding, errors=errors) as f:
            return f.read(n)
    read = read_text

    def readlines(self, *arg, **kw):
        return list(readrow(self, *arg, **kw))

    def groupbylines(self, *arg, **kw):
        return list(grouprow(self, *arg, **kw))

    def delete(self):
        return self.unlink()

    def lsdir(self, recursive=True):
        return [__class__(x) for x in lsdir(self, recursive)]

    def exists(self):
        try:
            return super().exists()
        except OSError:
            return False

    @property
    def encoding(self):
        return getencoding(self.read_bytes(SAMPLEBYTE))

    @property
    def ext(self):
        return self.info.ext

    @property
    def lineterminator(self):
        return getLF(self.read_bytes(SAMPLEBYTE))

    @property
    def owner(self):
        return self.info.owner

    @property
    def group(self):
        return self.info.group

    @property
    def permission(self):
        return self.info.permission

    @property
    def mtime(self):
        return self.info.mtime

    @property
    def ctime(self):
        return self.info.ctime

    @property
    def atime(self):
        return self.info.atime

    @property
    def size(self):
        return self.info.size

    def wordcount(self, word, buf_size = 1024 ** 2):
        with open(self.__str__(), mode="rb") as f:
            read_f = f.read # loop optimization
            if isinstance(word, str):
                if self.encoding:
                    word = word.encode(self.encoding)
                else:
                    word = word.encode()
            elif isinstance(word, (int, float)):
                word = bytes(str(word))

            pos = f.tell()
            if pos != 0:
                f.seek(0)

            buf = read_f(buf_size)

            words = 0
            while buf:
                words += buf.count(word)
                buf = read_f(buf_size)

            f.seek(pos)
            return words

    def linecount(self, buf_size = 1024 ** 2):
        return self.wordcount(word=b"\n", buf_size = buf_size)

    def uncompressedsize(self):
        size = getsize(self.__str__())
        if hasattr(size, "__next__"):
            return sum(s.size for s in size if s.target)
        return size

    def guesstype(self):
        if self._type is None:
            self._type = guesstype(self.__str__())
        return self._type

    def tree_file(self, recursive:bool=True, dotfile:bool=False):
        if dotfile:
            for r in self.lsdir(recursive):
                if r.is_file():
                    yield r
        else:
            for r in self.lsdir(recursive):
                if not r.is_file() or any(x.startswith(".") for x in r.parts):
                    continue
                yield r

    def tree_dir(self, recursive:bool=True, dotfile:bool=False):
        if dotfile:
            for r in self.lsdir(recursive):
                if r.is_dir():
                    yield r
        else:
            for r in self.lsdir(recursive):
                if not r.is_dir() or any(x.startswith(".") for x in r.parts):
                    continue
                yield r

    def is_compress(self):
        return self.guesstype() in ["zip", "gz", "bz2", "xz", "lha", "rar"]

    def open(self, *args, **kw):

        def pptx(path_or_buffer, *args, **kw):
            return Presentation(binopen(path_or_buffer), *args, **kw)

        def ppt(path_or_buffer, *args, **kw):
            with binopen(path_or_buffer) as fp:

                mgl.textdump = b""
                mgl.params.noStructOutput = True
                mgl.params.dumpText = True
                mgl.params.noRawDumps = True
                mgl.params.showSectorChain = True

                return PPTFile(fp.read(), mgl.params, *args, **kw)


        def docx(path_or_buffer, *args, **kw):
            return Document(binopen(path_or_buffer), *args, **kw)

        def doc(path_or_buffer, *args, **kw):
            return MSWordOLE(binopen(path_or_buffer), *args, **kw)

        def xlsx(path_or_buffer, *args, **kw):
            with binopen(path_or_buffer) as fp:
                return xlrd.open_workbook(file_contents=fp.read(), *args, **kw)

        xls = xlsx

        def pdf(path_or_buffer, *args, **kw):
            ps = PDFParser(binopen(path_or_buffer), *args, **kw)
            doc = PDFDocument()
            ps.set_document(doc)
            doc.set_parser(ps)
            doc.initialize('')
            return doc

        def csv(path_or_buffer, *args, **kw):
            with binopen(path_or_buffer) as fp:
                dat = fp.read(SAMPLEBYTE) + fp.readline()
                e = getencoding(dat)
                txt = dat.decode(e)
                dialect = _csv.Sniffer().sniff(txt)

            return _csv.reader(opener(path_or_buffer, encoding=e), dialect=dialect, *args, **kw)

        def dml(path_or_buffer, excludes = ["^;?\n", "^;$", "^\s*//.*$"], lvsep = "  "):
            with binopen(path_or_buffer) as fp:
                dat = fp.read()
            ret = []
            e = getencoding(dat)

            lines = list(enumerate((dat.decode(e) if e else dat.decode()).splitlines(),1))
            maxlen = max(x.count(lvsep) for i,x in lines)

            rec = [""] * maxlen

            for i, line in reversed(lines):

                if line and all(not re.match(x, line) for x in excludes):
                    row = line.rstrip().split(lvsep)
                    name = row[-1]
                    if name.startswith("end"):
                        recordname = name.replace("end", "").strip("[; ]")
                        rec[row.index(name)] = recordname
                    elif name.startswith("record"):
                        for x in range(row.index(name), maxlen):
                            rec[x] = ""
                    else:
                        ret.append([i, ".".join(rec).strip("."), line.strip()])

            for i, rec, line in reversed(ret):
                yield rec, line

        txt = opener

        def zip(path_or_buffer, *args, **kw):
            return ZipFile(binopen(path_or_buffer), *args, **kw)

        def gz(path_or_buffer, *args, **kw):
            return GzipFile(fileobj=binopen(path_or_buffer), *args, **kw)

        def bz2(path_or_buffer, *args, **kw):
            return BZ2File(binopen(path_or_buffer), *args, **kw)

        def xz(path_or_buffer, *args, **kw):
            return LZMAFile(binopen(path_or_buffer), *args, **kw)

        def tar(path_or_buffer, mode="r", *args, **kw):
            return TarFile.open(mode=mode, fileobj=binopen(path_or_buffer), *args, **kw)

        def lha(path_or_buffer, *args, **kw):
            return LhaFile(binopen(path_or_buffer), *args, **kw)

        def rar(path_or_buffer, *args, **kw):
            return RarFile(binopen(path_or_buffer), *args, **kw)

        def locate(path_or_buffer):
            fp = binopen(path_or_buffer)

            fp.read(8) # magic number
            conf_size = struct.unpack(">i", fp.read(4))[0]
            fp.read(4 + conf_size) # dbconf

            for entry in binchunk(fp, sep=b"\x00\x02"):
                raw = entry[16:].rstrip()
                raw = raw.replace(b"\x00\x00", b"\vfile\t").replace(b"\x00\x01", b"\vdir\t")

                root = ""
                for e in raw.decode("utf-8").split("\v"):
                    try:
                        tp, name = e.split("\t")
                        if name == "/":
                            tp = "dir"
                        else:
                            name = root + "/" + name
                            if name in ("/etc/init.d", "/etc/rc0.d", "/etc/rc1.d", "/etc/rc2.d", "/etc/rc3.d", "/etc/rc4.d", "/etc/rc5.d", "/etc/rc6.d", "/etc/ssl/certs", "/etc/xdg/systemd/user", "/lib", "/lib64", "/sbin", "/usr/lib64/go/4.8.5", "/usr/libexec/gcc/x86_64-redhat-linux/4.8.5", "/bin", "/usr/include/c++/4.8.5", "/usr/lib/debug/bin", "/usr/lib/debug/lib", "/usr/lib/debug/lib64", "/usr/lib/debug/sbin", "/usr/lib/gcc/x86_64-redhat-linux/4.8.5", "/usr/lib/go/4.8.5", "/usr/lib/terminfo", "/usr/share/doc/git-1.8.3.1/contrib/hooks", "/usr/share/doc/redhat-release", "/usr/share/doc/vim-common-7.4.160/docs", "/usr/share/gcc-4.8.5", "/usr/share/gccxml-0.9/GCC/5.0", "/usr/share/gccxml-0.9/GCC/5.1", "/usr/share/gccxml-0.9/GCC/5.2", "/usr/share/gccxml-0.9/GCC/5.3", "/usr/share/gdb/auto-load/lib64", "/usr/share/groff/current", "/usr/tmp", "/var/lock", "/var/mail", "/var/run"):
                                tp = "dir"

                        yield tp, name

                    except ValueError:
                        root = e

        def json(path_or_buffer, *args, **kw):
            with binopen(path_or_buffer) as fp:
                return jsonlib.loads(fp.read(), *args, **kw)

        def html(path_or_buffer, *args, **kw):
            with binopen(path_or_buffer) as fp:
                dat = fp.read()
            return _HTMLParser(dat, getencoding(dat), *args, **kw)

        def xml(path_or_buffer, *args, **kw):
            with binopen(path_or_buffer) as fp:
                dat = fp.read()
            return xmltodict(dat, getencoding(dat), *args, **kw)

        def mdb(path_or_buffer, uid="", passwd="", *args, **kw):
            with binopen(path_or_buffer) as fp:
                path = os.path.abspath(fp.name)
            try:
                driver = next(x for x in pyodbc.drivers() if x.startswith("Microsoft Access Driver "))
            except StopIteration:
                raise RuntimeError("Not Installed Microsoft Access Driver")

            dsnstr = r'DRIVER={{{}}};DBQ={};UID="{}";PWD="{}";'
            return pyodbc.connect(dsnstr.format(driver, path, uid, passwd), *args, **kw)

        accdb = mdb

        def sqlite3(path_or_buffer, *args, **kw):
            with binopen(path_or_buffer) as fp:
                path = os.path.abspath(fp.name)
            return create_engine(r"sqlite:///{}".format(path), *args, **kw)

        def pickle(path_or_buffer, *args, **kw):
            with binopen(path_or_buffer) as fp:
                return pkload(fp, *args, **kw)

        return locals().get(guesstype(self.__str__()), open)(self.__str__(), *args, **kw)


    def close(self):
        if hasattr(self._file, "close") and self._file.closed is False:
            self._file.close()
            self._file = None
            self._rows = None

    def __next__(self):
        return next(self.__iter__())

    def __iter__(self):
        if self._rows is None:
            self._rows = readrow(self)
        return self._rows

    def __enter__(self):
        if self._closed:
            self._raise_closed()
        return self.open()

    def __exit__(self, t, v, tb):
        self.close()


def test():
    from util.core import tdir
    from glob import glob
    from datetime import datetime as dt

    def test_lsdir():
        def _testargs(func, pathstr, *args):
            #TODO assert
            ret = [
                func(pathstr, *args),
                func(pathlib.Path(pathstr), *args),
                func(pathlib.Path(pathstr), *args),
            ]
            try:
                with open(pathstr) as f:
                    ret.append(func(f, *args))
            except Exception as e:
                ret.append(e)
            return ret
        _testargs(lsdir, tdir)
        _testargs(lsdir, tdir + "diff*")
        _testargs(lsdir, tdir+"test.csv")
        _testargs(lsdir, tdir+"*est.csv")
        _testargs(lsdir, tdir+"ddfghjdtui")
        _testargs(lsdir, tdir, False)
        _testargs(lsdir, tdir + "diff*", False)
        _testargs(lsdir, tdir+"test.csv", False)
        _testargs(lsdir, tdir+"*est.csv", False)
        _testargs(lsdir, tdir+"ddfghjdtui", False)

    def test_readrow():
        for g in glob(tdir+"*.*"):
            n = os.path.basename(g)

            sys.stdout.write("test_readrow : " + n)
            try:
                assert list(readrow(g))
                print(" ok ")
            except (NotImplementedError):
                print(" skip")
            except TypeError:
                print(" None")
            except Exception as e:
                print(" ERROR", e)

    def test_grouprow():
        for g in glob(tdir+"*.*"):
            n = os.path.basename(g)

            sys.stdout.write("test_grouprow : " + n)
            try:
                assert list(grouprow(g))
                print(" ok ")
            except (NotImplementedError):
                print(" skip")
            except TypeError:
                print(" None")
            except Exception as e:
                print(" ERROR", e)

    def test_info():
        for g in glob(tdir+"*.*"):
            n = os.path.basename(g)

            sys.stdout.write("test_info : " + n)
            try:
                assert list(getinfo(g))
                print(" ok ")
            except (NotImplementedError):
                print(" skip")
            except TypeError:
                print(" None")
            except Exception as e:
                print(" ERROR", e)

    def test_getsize():
        for g in glob(tdir+"*.*"):
            n = os.path.basename(g)

            sys.stdout.write("test_getsize : " + n)
            try:
                print(" ok ", getsize(g))
            except (NotImplementedError):
                print(" skip")
            except TypeError:
                print(" None")
            except Exception as e:
                print(" ERROR", e)

    def test_Path_init():
        for g in glob(tdir+"*.*"):
            n = os.path.basename(g)

            sys.stdout.write("test_Path : " + n)
            try:
                p = Path(g)
                assert p
                assert p.size > 0
                assert "LazyInfo" in str(type(p.info))
                assert type(p.encoding) in [str, type(None)]
                p.guesstype()
                dat = p.read_bytes(SAMPLEBYTE)
                if b"\x00" not in dat:
                    assert p.read_text(10)
                assert p.readlines()
                assert p.groupbylines()

                assert(type(p.is_compress()) is bool)
                assert p.uncompressedsize() > 0
                assert next(p)
                p._rows = None
                assert list(p)
                assert p.open()
                with p as f:
                    assert f
                p.close()

                print(" ok ")
            except NotImplementedError as e:
                print(" skip", e)
            except TypeError as e:
                print(" None", e)
            except Exception as e:
                print(" ",e.__class__.__name__, e)


    t0 = dt.now()
    for x, func in list(locals().items()):
        if x.startswith("test_") and callable(func):
            t1 = dt.now()
            func()
            t2 = dt.now()
            print("{} : time {}".format(x, t2-t1))
    t3 = dt.now()
    if x.startswith("test_"):
        print("{} : time {}".format(x, t3-t0))

from glob import glob
def walk(args):
    for arg in args.files:
        for f in glob(arg):
            if os.path.isdir(f):
                continue
            f = os.path.normpath(f)

            if args.verbose:
                sys.stderr.write("Dumping:{}\n".format(f))
                sys.stderr.flush()
            yield f

def unicode_escape(x):
    return x.encode().decode("unicode_escape")

def main_row():
    from argparse import ArgumentParser

    ps = ArgumentParser(prog="dumper",
                        description="any file text dump\n")
    padd = ps.add_argument

    padd('-V','--version', action='version', version='%(prog)s ' + __version__)
    padd("-v", "--verbose", help="print progress",
         action='store_true', default=False)
    padd("files",
         metavar="<files>",
         nargs="+",  default=[],
         help="text dump any files")

    padd('-e', '--encoding', type=str, default="cp932",
         help='output fileencoding (default `cp932`)')
    padd('-s', '--sep', type=unicode_escape, default="\t",
         help='output separater (default `\\t`)')
    padd('-l', '--lineterminator', type=unicode_escape, default="\r\n",
         help='output fileencoding (default `\\r\\n`)')
    padd('-f', '--filename', action='store_true', default=False,
         help='output filename (default False)')
    padd('-t', '--target', action='store_true', default=False,
         help='output targetname (default False)')
    args = ps.parse_args()

    encoding = args.encoding
    lineterminator = args.lineterminator
    sep = args.sep
    target = args.target
    filename = args.filename

    i = None
    write = sys.stdout.buffer.write

    def oneliner(b:str):
        if b is None:
            return ""
        try:
            if "\r" in b:
                b = b.replace("\r", "\\r")
            if "\n" in b:
                b = b.replace("\n", "\\n")
            return b
        except TypeError:
            return str(b)

    for i, f in enumerate(walk(args)):
        for x in readrow(f):
            if filename:
                write("{}:\t".format(x.path).encode(encoding))
            if target:
                write("[{}]\t".format(x.target).encode(encoding))
            try:
                write((x.value + lineterminator).encode(encoding))
            except TypeError:
                write((sep.join(map(oneliner,x.value)) + lineterminator).encode(encoding))

    if i is None:
        raise FileNotFoundError(str(args.files))


def main_info():
    from argparse import ArgumentParser

    ps = ArgumentParser(prog="info",
                        description="any file infomation program\n")
    padd = ps.add_argument

    padd('-V','--version', action='version', version='%(prog)s ' + __version__)
    padd("-v", "--verbose", help="print progress",
         action='store_true', default=False)
    padd("files",
         metavar="<files>",
         nargs="+",  default=[],
         help="infomation filenames")

    padd('-s', '--sep', type=unicode_escape, default=",",
         help='output separater (default `,`)')
    padd('-l', '--lineterminator', type=unicode_escape, default="\r\n",
         help='output fileencoding (default `\\r\\n`)')
    padd('-d', '--dateformat', type=str, default="%Y/%m/%d %H:%M:%S",
         help='output datetimeformat (default `%Y/%m/%d %H:%M:%S`)')

    args = ps.parse_args()

    sep = args.sep
    lineterminator = args.lineterminator
    dateformat = args.dateformat

    i = None
    write = sys.stdout.write

    write((sep.join(sinfo._fields) + lineterminator))

    for i, f in enumerate(walk(args)):
        ret = [x if x else "" for x in getinfo(f)]
        idx = sinfo._fields.index("LF")
        ret[idx] = repr(ret[idx])
        for j, r in enumerate(ret):
            if isinstance(r, datetime):
                ret[j] = r.strftime(dateformat)
        write((sep.join(map(str, ret)) + lineterminator))

    if i is None:
        raise FileNotFoundError(str(args.files))

def main_size():
    from argparse import ArgumentParser

    ps = ArgumentParser(prog="info",
                        description="any file infomation program\n")
    padd = ps.add_argument

    padd('-V','--version', action='version', version='%(prog)s ' + __version__)
    padd("-v", "--verbose", help="print progress",
         action='store_true', default=False)
    padd("files",
         metavar="<files>",
         nargs="+",  default=[],
         help="infomation filenames")

    padd('-u', '--unit', type=str, default="K",
         help='output file unit byte (default `K`) choice=> (`B`, `K`, `M`, `G`)')

    padd('-l', '--lineterminator', type=unicode_escape, default="\r\n",
         help='output fileencoding (default `\\r\\n`)')

    args = ps.parse_args()

    lineterminator = args.lineterminator
    unitname = args.unit.upper()
    unit = dict(B=1, K=1024, M=1024**2, G=1024**3)[unitname]
    if unitname == "B":
        form = "{}: {:.0f} " + "Byte" + lineterminator
    else:
        form = "{}: {:.2f} " + unitname + "Byte" + lineterminator

    i = None
    write = sys.stdout.write

    for i, f in enumerate(walk(args)):
        write(form.format(f, getsize(f) / unit))

    if i is None:
        raise FileNotFoundError(str(args.files))

def main():
    try:
        subcmd = sys.argv.pop(1)
    except IndexError:
        sys.stderr.write("Not Found subcommand.\nPlease input subcommand(`row` or `info` or `size`)")
        sys.exit(1)
    if subcmd in ["row", "info", "size"]:
        eval("main_" + subcmd)()
    else:
        raise AttributeError("python {} [row|info|size] ...".format(os.path.basename(sys.argv[0])))

if __name__ == "__main__":
#    test()
    main()
